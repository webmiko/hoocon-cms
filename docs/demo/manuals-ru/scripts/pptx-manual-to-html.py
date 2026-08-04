#!/usr/bin/env python3
"""Build RU landscape A4 HTML manuals (2 sheets) for Hoocon actuators.

Canon for specs/diagrams: ``_инструкции-pdf/*.pdf`` (via DiagramProfile +
normalizers). PPTX under Downloads is only a layout/text scaffold for
extraction — not published.

Voltage shells for EN→RU (separate 24 V / 230 V PDFs): ``VoltageTemplate`` /
``VOLTAGE_TEMPLATES`` → ``template-v24.html`` / ``template-v230.html``.

Glossary: docs/tech-copy-belimo-ru.md (Belimo RU canon).
Layout canon: docs/demo/manuals-ru/TEMPLATES.md
  (12-col A4×2; wiring → dimensions → rotation; V24/V230 vs combined).

Run from repo root::

    python3 docs/demo/manuals-ru/scripts/pptx-manual-to-html.py
"""

from __future__ import annotations

import argparse
import html
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

EMU_PER_MM = 914400 / 25.4

# This file lives at ``docs/demo/manuals-ru/scripts/…``.
_MANUALS_RU_DIR = Path(__file__).resolve().parents[1]
# Published finished manuals by product family (EN→RU queue).
MANUAL_FAMILY_DIRS: dict[str, str] = {
    "da": "DA",
    "sa": "SA",
    "hv": "HV",
}
FINISHED_MANUALS_SUBDIR = "DA"  # default / legacy alias for DA family


def manual_family_for_stem(stem: str) -> str:
    """Map HTML stem to family folder key: da | sa | hv."""
    s = stem.casefold()
    if s.startswith("sa"):
        return "sa"
    if s.startswith("hv"):
        return "hv"
    return "da"


def finished_manuals_subdir(stem: str | None = None, *, family: str | None = None) -> str:
    """Directory name under manuals-ru for finished HTML (``DA`` / ``SA`` / ``HV``)."""
    key = family or (manual_family_for_stem(stem) if stem else "da")
    return MANUAL_FAMILY_DIRS.get(key, FINISHED_MANUALS_SUBDIR)


def finished_manuals_dir(
    out_dir: Path | None = None,
    *,
    stem: str | None = None,
    family: str | None = None,
) -> Path:
    """Published finished manuals live under ``manuals-ru/{DA,SA,HV}/``."""
    base = out_dir if out_dir is not None else _MANUALS_RU_DIR
    return base / finished_manuals_subdir(stem, family=family)


def _ensure_family_logo(finished_dir: Path) -> None:
    """Symlink Hoocon logo into ``{DA,SA,HV}/assets/`` for relative HTML paths."""
    assets = finished_dir / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    link = assets / "hoocon-logo.svg"
    target = Path("../../assets/hoocon-logo.svg")
    if link.is_symlink() or link.is_file():
        return
    link.symlink_to(target)


def iter_finished_manual_html(
    out_dir: Path | None = None,
) -> list[tuple[str, str, Path]]:
    """Yield ``(family_subdir, stem, html_path)`` for all finished manuals."""
    base = out_dir if out_dir is not None else _MANUALS_RU_DIR
    rows: list[tuple[str, str, Path]] = []
    for fam_dir in MANUAL_FAMILY_DIRS.values():
        folder = base / fam_dir
        if not folder.is_dir():
            continue
        for html_path in sorted(folder.glob("*.html")):
            rows.append((fam_dir, html_path.stem, html_path))
    return rows


def _repo_root() -> Path:
    """Hoocon-cms repo root (has ``backend/`` and ``_инструкции-pdf/``)."""
    for parent in Path(__file__).resolve().parents:
        if (parent / "backend").is_dir() and (parent / "_инструкции-pdf").is_dir():
            return parent
    return Path(__file__).resolve().parents[4]


# Longest phrases first. Canon: docs/tech-copy-belimo-ru.md
PHRASES: list[tuple[str, str]] = [
    (
        "1. The Damper Actuator is not allowed to be used outside the specified "
        "field of application, especially in aircraft.\n2. The enclosure of the "
        "actuator equipment may only be opened by the manufacturer. It contains "
        "no components which the user can replace or repair.\n\n3. The device "
        "contains electrical and electronic components and is not allowed to be "
        "disposed of as household refuse. All locally valid regulations and "
        "requirements must be observed.",
        "1. Запрещается использовать электропривод заслонки вне указанной области "
        "применения, особенно в авиационной технике.\n2. Вскрытие корпуса привода "
        "разрешено только производителю. Внутри нет компонентов, которые "
        "пользователь может заменять или ремонтировать.\n\n3. Устройство содержит "
        "электрические и электронные компоненты, в связи с чем недопустима "
        "утилизация вместе с бытовыми отходами. Необходимо соблюдать все "
        "действующие правила и инструкции, относящиеся к данной конкретной местности.",
    ),
    (
        "1.Do not use damper actuators on not specified application range, "
        "especially do not use damper actuators on aircraft.\n"
        "2.Nonprofessional installation personnel are not allowed to open casing "
        "of damper actuator, and do not connect cables when power is on.\n"
        "3.The damper actuators contain electronic components, which shall not "
        "treated as ordinary household waste, and shall be disposed in accordance "
        "with the relevant local laws and regulations.",
        "1. Запрещается использовать электроприводы заслонок вне указанной области "
        "применения, особенно в авиационной технике.\n"
        "2. Неспециалистам запрещается вскрывать корпус привода; подключение "
        "кабелей выполнять только при отключённом питании.\n"
        "3. Приводы содержат электронные компоненты и не подлежат утилизации вместе "
        "с бытовыми отходами; соблюдайте местные нормы утилизации.",
    ),
    (
        "For operation of air control dampers in HVAC system",
        "Для управления воздушными заслонками в системах HVAC "
        "(отопления, вентиляции и кондиционирования воздуха)",
    ),
    (
        "SPRING-RETURN DAMPER ACTUATOR\nInstruction Manual",
        "Привод с пружинным возвратом\nРуководство по эксплуатации",
    ),
    (
        "SPRING-RETURN DAMPER ACTUATOR Instruction Manual",
        "Привод с пружинным возвратом\nРуководство по эксплуатации",
    ),
    (
        "FIRE AND SMOKE DAMPER ACTUATOR\nInstruction Manual",
        "Привод противопожарного клапана\nРуководство по эксплуатации",
    ),
    (
        "FIRE AND SMOKE DAMPER ACTUATOR Instruction Manual",
        "Привод противопожарного клапана\nРуководство по эксплуатации",
    ),
    (
        "FIRE AND SMOKE DAMPER ACTUATOR",
        "Привод противопожарного клапана",
    ),
    (
        "GENERAL DAMPER ACTUATOR\nInstruction Manual",
        "Привод воздушной заслонки\nРуководство по эксплуатации",
    ),
    (
        "GENERAL DAMPER ACTUATOR Instruction Manual",
        "Привод воздушной заслонки\nРуководство по эксплуатации",
    ),
    (
        "Max 45dB(A)when motor working, max 62dB(A) while spring return",
        "макс. 45 дБ(А) при работе двигателя, макс. 62 дБ(А) при возврате пружины",
    ),
    (
        "Max 45dB(A)when motor working, max 50dB(A) while spring return",
        "макс. 45 дБ(А) при работе двигателя, макс. 50 дБ(А) при возврате пружины",
    ),
    (
        "The direction of rotation can be changed by the screw on top cover",
        "Направление вращения можно изменить винтом на верхней крышке",
    ),
    (
        "The direction of rotation can be changed by the motor pin.",
        "Направление вращения можно изменить с помощью перемычки на двигателе.",
    ),
    (
        "Mechanical position limitation installation method",
        "Установка механических ограничителей угла поворота",
    ),
    (
        "12x12mm □ Can provide 8x8,10x10mm shaft sleeve",
        "Квадратный 12×12 мм (доступны втулки 8×8, 10×10 мм)",
    ),
    (
        "Max.95°, adjustable by mechanical stops",
        "Макс. 95°, регулируется механическими упорами",
    ),
    (
        "Max.95°,adjustable by mechanical stops",
        "Макс. 95°, регулируется механическими упорами",
    ),
    (
        "Protection class Ⅲ (safety low voltage)",
        "Класс защиты III (безопасное сверхнизкое напряжение)",
    ),
    (
        "Ⅲ (safety low voltage)",
        "III (безопасное сверхнизкое напряжение)",
    ),
    (
        "Ⅱ (totally insulated)",
        "II (все изолировано / полная изоляция)",
    ),
    (
        "95%RH, non condensing／EN 60730-1",
        "95% относительной влажности, без конденсации / согласно EN 60730-1",
    ),
    (
        "Dial switch set of control signal",
        "Поворотный переключатель режима управляющего сигнала",
    ),
    (
        "Switch The Direction Of Rotation",
        "Переключение направления вращения",
    ),
    (
        "Switch direction of rotation",
        "Переключение направления вращения",
    ),
    (
        "to install on the other side",
        "для монтажа с противоположной стороны",
    ),
    (
        "Through manual button",
        "Кнопка ручного управления",
    ),
    (
        "Through metal handle",
        "Металлическая рукоятка",
    ),
    (
        "Selectable by switch",
        "Выбирается с помощью переключателя",
    ),
    (
        "AS type include 2 groups of auxiliary switches.",
        "Исполнение «AS» включает 2 группы вспомогательных переключателей.",
    ),
    (
        "AS type include 2 groups of auxiliary switches",
        "Исполнение «AS» включает 2 группы вспомогательных переключателей",
    ),
    (
        "DS type include 1 group of auxiliary switch",
        "Исполнение «DS» включает 1 группу вспомогательных переключателей",
    ),
    (
        "S type include 2 auxiliary switch.",
        "Исполнение «S» включает 2 вспомогательных переключателя.",
    ),
    (
        "S type include 2 auxiliary switch",
        "Исполнение «S» включает 2 вспомогательных переключателя",
    ),
    (
        "Control: 0(2)...10VDC/0(4)...20mA",
        "Упр. сигнал Y: 0(2)...10 В= / 0(4)...20 мА (спецзаказ)",
    ),
    (
        "Feedback: 0(2)...10VDC/0(4)...20mA",
        "Обратная связь U: 0(2)...10 В= / 0(4)...20 мА (спецзаказ)",
    ),
    (
        "Control: ON/OFF, 2/3-Point",
        "Управление: вкл./выкл. (ON/OFF), 2- или 3-позиционное",
    ),
    (
        "Control: ON/OFF",
        "Управление: вкл./выкл. (ON/OFF)",
    ),
    (
        "10/15/20Nm Modulating",
        "10/15/20 Нм, пропорциональное (модулирующее) управление",
    ),
    (
        "10/15/20Nm ON/OFF",
        "10/15/20 Нм, вкл./выкл. (ON/OFF)",
    ),
    (
        "4/6Nm Modulating",
        "4/6 Нм, пропорциональное (модулирующее) управление",
    ),
    (
        "4Nm/6Nm ON/OFF",
        "4/6 Нм, вкл./выкл. (ON/OFF)",
    ),
    (
        "2Nm Modulating",
        "2 Нм, пропорциональное (модулирующее) управление",
    ),
    (
        "2Nm ON/OFF",
        "2 Нм, вкл./выкл. (ON/OFF)",
    ),
    (
        "3Nm ON/OFF",
        "3 Нм, вкл./выкл. (ON/OFF)",
    ),
    (
        "Running time：  <100s, 150s, 170s Spring reset time:<25s",
        "Время поворота: < 100 с / 150 с / 170 с; время возврата пружины: < 25 с",
    ),
    (
        "Running time：  <110s, Spring reset time:<25s",
        "Время поворота: < 110 с; время возврата пружины: < 25 с",
    ),
    (
        "Running time：  <110s, 150s, Spring reset time:<25s",
        "Время поворота: < 110 с / 150 с; время возврата пружины: < 25 с",
    ),
    (
        "Running time：  <110s/150s, Spring reset time:<25s",
        "Время поворота: < 110 с / 150 с; время возврата пружины: < 25 с",
    ),
    (
        "Running time：  <75s, Spring reset time:<25s",
        "Время поворота: < 75 с; время возврата пружины: < 25 с",
    ),
    (
        "Running time：<30s   (95°)",
        "Время поворота: < 30 с (на угол 95°)",
    ),
    (
        "Torque: 10Nm, 15Nm, 20Nm",
        "Крутящий момент: 10 Нм, 15 Нм, 20 Нм",
    ),
    (
        "Torque: 4Nm/6Nm",
        "Крутящий момент: 4 Нм / 6 Нм",
    ),
    (
        "Torque: 2Nm",
        "Крутящий момент: 2 Нм",
    ),
    (
        "Torque: 3Nm",
        "Крутящий момент: 3 Нм",
    ),
    (
        "Nominal voltage：AC/DC 24V   AC 100~240V",
        "Номинальное напряжение: AC/DC 24 В; AC 100…240 В",
    ),
    (
        "Nominal voltage：AC/DC 24V AC100~240V",
        "Номинальное напряжение: AC/DC 24 В; AC 100…240 В",
    ),
    (
        "Nominal voltage ：AC/DC 24V   AC 100~240V",
        "Номинальное напряжение: AC/DC 24 В; AC 100…240 В",
    ),
    (
        "Nominal voltage：AC/DC 24V",
        "Номинальное напряжение: AC/DC 24 В",
    ),
    (
        "10W@ Nominal toque/3W@holding",
        "10 Вт под нагрузкой\n3 Вт удержание",
    ),
    (
        "10W@Nominal toque/0.5W@holding",
        "10 Вт под нагрузкой\n0,5 Вт удержание",
    ),
    (
        "6W@ Nominal toque/1.5W@holding",
        "6 Вт под нагрузкой\n1,5 Вт удержание",
    ),
    (
        "6W@Nominal toque/1.5W@holding",
        "6 Вт под нагрузкой\n1,5 Вт удержание",
    ),
    (
        "5W@ Nominal toque/3W@holding",
        "5 Вт под нагрузкой\n3 Вт удержание",
    ),
    (
        "5W@ Nominal toque/2W@holding",
        "5 Вт под нагрузкой\n2 Вт удержание",
    ),
    (
        "3W@Nominal toque/0.8W@holding",
        "3 Вт под нагрузкой\n0,8 Вт удержание",
    ),
    (
        "3W@Nominal toque/0.7W@holding",
        "3 Вт под нагрузкой\n0,7 Вт удержание",
    ),
    (
        "3W@Nominal toque/0.5W@holding",
        "3 Вт под нагрузкой\n0,5 Вт удержание",
    ),
    (
        "NINGBO HOOCON AUTOMATION CONTROL EQUIPMENTCO.,LTD.",
        "NINGBO HOOCON AUTOMATION CONTROL EQUIPMENT CO., LTD.",
    ),
    (
        "NINGBOHOOCONAUTOMATIONCONTROLEQUIPMENTCO.,LTD.",
        "NINGBO HOOCON AUTOMATION CONTROL EQUIPMENT CO., LTD.",
    ),
    (
        "Add:No.1 Licui Road Chaotang Industrial Zone Zonghan,Cixi,China.315301",
        "Адрес: No.1 Licui Road, Chaotang Industrial Zone, Zonghan, Cixi, China. 315301",
    ),
    (
        "Add:No.298ChaotangIndustrialZoneZonghan,Cixi,China.315301",
        "Адрес: No.298 Chaotang Industrial Zone, Zonghan, Cixi, China. 315301",
    ),
    ("Tel:+86-574-63813330", "Тел.: +86-574-63813330"),
    ("Fax:+86-574-63220759", "Факс: +86-574-63220759"),
    ("E-mail:hoocon@hoocon.com.cn", "E-mail: hoocon@hoocon.com.cn"),
    ("Http://www.hoocon.com.cn", "http://www.hoocon.com.cn"),
    ("Technical specification", "Технические характеристики"),
    ("Technical specifi cation", "Технические характеристики"),
    ("Actuator Dimensions(mm)", "Габаритные размеры привода (мм)"),
    ("Wiring Diagram", "Схема подключения"),
    ("Auxiliary switch", "Вспомогательный переключатель"),
    ("Electrical data", "Электрические параметры"),
    ("Function data", "Функциональные параметры"),
    ("Working/conditions", "Условия эксплуатации"),
    ("Dimensions/Weight", "Габаритные размеры / Масса"),
    ("Dimensions(L ×W ×H)", "Габаритные размеры\n(Д × Ш × В)"),
    ('See"Dimensions"', "См. «Габаритные чертежи»"),
    ("Nominal voltage", "Номинальное напряжение"),
    ("Power consumption", "Потребляемая мощность"),
    ("Wire sizing", "Сечение подключаемых проводов"),
    ("Damper size", "Площадь обслуживаемой заслонки"),
    ("Direction of rotation", "Направление вращения"),
    ("Manual override", "Ручное управление"),
    ("Angle of rotation", "Угол поворота"),
    ("Sound power level", "Уровень звуковой мощности"),
    ("Position indication", "Индикация положения"),
    ("Protection class", "Класс защиты"),
    ("Protection level", "Степень защиты корпуса"),
    ("Ambient temperature", "Температура окружающей среды"),
    ("Inventory temperature", "Температура хранения"),
    ("Humidity test", "Испытание на влажность"),
    ("Shaft length", "Длина вала заслонки"),
    ("Shaft diameter", "Диаметр вала заслонки"),
    ("Running time", "Время поворота"),
    ("Rotary switch", "Поворотный переключатель"),
    ("Short circuit", "Замкнуто"),
    ("Open circuit", "Разомкнуто"),
    ("Switch a", "Переключатель a"),
    ("Switch b", "Переключатель b"),
    ("Attention:", "ВНИМАНИЕ:"),
    ("Actuator", "Привод"),
    ("Weight", "Масса"),
    ("Mechanical", "Механический указатель"),
    ("Max95°", "Макс. 95°"),
    ("Max.95°", "Макс. 95°"),
    ("AC/DC 24V 50/60Hz", "AC/DC 24 В, 50/60 Гц"),
    ("AC100-240V 50/60Hz", "AC 100…240 В, 50/60 Гц"),
    ("DC24V AC230V", "DC 24 В / AC 230 В"),
    ("AC24V", "AC 24 В"),
    ("Terminal 21，22", "Клеммы 21, 22"),
    ("Terminal 21，23", "Клеммы 21, 23"),
    ("Terminal21，22", "Клеммы 21, 22"),
    ("Terminal21，23", "Клеммы 21, 23"),
    ("Terminal24，25", "Клеммы 24, 25"),
    ("Terminal24，26", "Клеммы 24, 26"),
    ("<70s(95°)", "< 70 с (95°)"),
    ("<50s(95°)", "< 50 с (95°)"),
    ("<100s/ <25s", "< 100 с / < 25 с"),
    ("<150s/ <25s", "< 150 с / < 25 с"),
    ("<170s/ <25s", "< 170 с / < 25 с"),
    ("-20... +50℃", "–20…+50 °C"),
    ("-40... +70℃", "–40…+70 °C"),
    ("-30... +80℃", "–30…+80 °C"),
    ("-30.. +80℃", "–30…+80 °C"),
    ("5~95%RH", "5…95 % отн. влажности"),
    ("○ 6...16mm □ 5x5...12x12mm", "круглый 6…16 мм, квадратный 5×5…12×12 мм"),
    ("○ 6...16mm □ 8x8...12x12mm", "круглый 6…16 мм, квадратный 8×8…12×12 мм"),
    ("○ 10...16mm □7...11mm", "круглый 10…16 мм, квадратный 7×7…11×11 мм"),
    ("○8...21mm □6...15mm", "круглый 8…21 мм, квадратный 6…15 мм"),
    ("＞ 50mm", "> 50 мм"),
    ("＞50mm", "> 50 мм"),
    ("＞ 90mm", "> 90 мм"),
    ("＜ 0.5kg", "< 0,5 кг"),
    ("＜ 0.7kg", "< 0,7 кг"),
    ("＜1.5kg", "< 1,5 кг"),
    ("＜1.3kg", "< 1,3 кг"),
    ("＜2.6kg", "< 2,6 кг"),
    ("＜ 2.50 kgs", "< 2,5 кг"),
    ("＜0.2m", "< 0,2 м²"),
    ("＜0.4 m", "< 0,4 м²"),
    ("＜0.6m", "< 0,6 м²"),
    ("＜0.2 2", "< 0,2 м²"),
    ("＜0.4 2", "< 0,4 м²"),
    ("＜0.6 2", "< 0,6 м²"),
    ("0. 3 m2", "0,3 м²"),
    ("0. 5 m2", "0,5 м²"),
    ("0.5m 2", "0,5 м²"),
    ("1.0 ㎡", "1,0 м²"),
    ("1.5 ㎡", "1,5 м²"),
    ("2.0 ㎡", "2,0 м²"),
    ("2.0㎡", "2,0 м²"),
    ("0.5mm", "0,5 мм²"),
    ("0.5m㎡", "0,5 мм²"),
    ("4Nm/6Nm", "4 Нм / 6 Нм"),
    ("10Nm/15Nm/20Nm", "10 / 15 / 20 Нм"),
    ("45dB", "45 дБ"),
    ("2Nm", "2 Нм"),
    ("( N o4 : U ndefined .)", "(№4: не используется)"),
    ("0-5°", "0–5°"),
    ("0-10°", "0–10°"),
    ("0-80°", "0–80°"),
    ("10-90°", "10–90°"),
    ("80-90°", "80–90°"),
]

_NBSP = re.compile(r"\u00a0")
_MULTI_WS = re.compile(r"[ \t\u00a0]+")
_BLANK_LINES = re.compile(r"\n{2,}")

# Label-only keys (safe mid-string only when the whole blob is short).
_LABELS = {
    "Technical specification",
    "Technical specifi cation",
    "Actuator Dimensions(mm)",
    "Wiring Diagram",
    "Auxiliary switch",
    "Electrical data",
    "Function data",
    "Working/conditions",
    "Dimensions/Weight",
    "Dimensions(L ×W ×H)",
    'See"Dimensions"',
    "Nominal voltage",
    "Power consumption",
    "Wire sizing",
    "Damper size",
    "Direction of rotation",
    "Manual override",
    "Angle of rotation",
    "Sound power level",
    "Position indication",
    "Protection class",
    "Protection level",
    "Ambient temperature",
    "Inventory temperature",
    "Humidity test",
    "Shaft length",
    "Shaft diameter",
    "Running time",
    "Rotary switch",
    "Short circuit",
    "Open circuit",
    "Switch a",
    "Switch b",
    "Attention:",
    "Actuator",
    "Weight",
    "Mechanical",
}


def _flex_pattern(eng: str) -> re.Pattern[str]:
    """Match English phrase allowing messy PPTX whitespace/newlines."""
    parts = [re.escape(p) for p in re.split(r"\s+", eng.strip()) if p]
    return re.compile(r"\s+".join(parts), re.IGNORECASE)


_FLEX: list[tuple[re.Pattern[str], str, str]] = [
    (_flex_pattern(eng), eng, rus) for eng, rus in PHRASES
]


def translate(text: str) -> str:
    """Apply glossary phrase map; keep SKU codes / numbers intact."""
    if not text or not text.strip():
        return text
    out = _NBSP.sub(" ", text)
    out = out.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse runs of spaces inside lines but keep newlines.
    out = "\n".join(_MULTI_WS.sub(" ", line).strip() for line in out.split("\n"))
    out = _BLANK_LINES.sub("\n\n", out).strip("\n")

    for pat, eng, rus in _FLEX:
        if eng in _LABELS and len(out) > 80:
            # Avoid eating "Actuator" inside long English attention paragraphs.
            continue
        out = pat.sub(rus, out)

    # Extra product-specific leftovers.
    extras = [
        (r"(?i)Spring reset time\s*:\s*", "время возврата пружины: "),
        (r"(?i)\bTorque:\s*", "Крутящий момент: "),
        (r"(?i)GENERAL\s+DAMPER\s+ACTUATOR", "Привод воздушной заслонки"),
        (r"(?i)Instruction\s+Manual", "Руководство по эксплуатации"),
        (r"(?i)FIRE\s+AND\s+SMOKE\s+DAMPER\s+ACTUATOR",
         "Привод противопожарного клапана"),
        (r"(?i)SPRING-RETURN\s+DAMPER\s+ACTUATOR",
         "Привод с пружинным возвратом"),
    ]
    for pat, rus in extras:
        out = re.sub(pat, rus, out)

    out = out.replace("VDC", "В=")
    out = re.sub(r"(?<![А-Яа-я])mA\b", "мА", out)
    out = re.sub(r"(?<=\d)\s*Nm\b", " Нм", out)
    out = out.replace("N·m", "Нм")
    out = re.sub(r"(?<=\d)V\b", " В", out)
    out = re.sub(r"(?<=\d)W\b", " Вт", out)
    out = re.sub(r"(?<=\d)s\b", " с", out)
    return out

from dataclasses import dataclass, field

_SKU_TOKEN = re.compile(
    r"\b(?:DA|SA|HV[AD]?|H8\d{0,3})[A-Z0-9]*(?:-\w+)?\b",
    re.IGNORECASE,
)


def iter_shapes(shapes, ox: int = 0, oy: int = 0):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh.shapes, ox + int(sh.left), oy + int(sh.top))
        else:
            yield sh, ox + int(sh.left or 0), oy + int(sh.top or 0)


def _shape_text(sh) -> str:
    if not sh.has_text_frame:
        return ""
    return "\n".join(p.text for p in sh.text_frame.paragraphs).strip()


def _table_matrix(sh) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in sh.table.rows:
        rows.append([translate(c.text.strip()) for c in row.cells])
    return rows


def _is_banner_row(matrix: list[list[str]]) -> bool:
    return len(matrix) == 1 and len(matrix[0]) == 1 and bool(matrix[0][0].strip())


def _is_dimensions_weight_table(matrix: list[list[str]]) -> bool:
    """Separate PPTX table for shaft / mass / L×W×H (not folded into function)."""
    if not matrix or len(matrix) > 8:
        return False
    joined = " ".join(c for row in matrix for c in row).lower()
    hits = sum(
        1
        for k in ("габарит", "масса", "вала", "dimensions", "weight", "shaft")
        if k in joined
    )
    return hits >= 2


def _function_has_dimensions(matrix: list[list[str]]) -> bool:
    for row in matrix:
        if not row:
            continue
        lab = row[0].strip().lower()
        if "габарит" in lab or lab == "масса" or "вала" in lab:
            return True
    return False


def _merge_dimensions_into_function(
    function: list[list[str]],
    dims: list[list[str]],
) -> list[list[str]]:
    """Append a split dimensions/weight table into the function matrix."""
    if not dims or not function or _function_has_dimensions(function):
        return function
    width = max(len(r) for r in function)
    out = [list(r) for r in function]
    # Section banner if missing.
    if not any(r and r[0].strip() == "Габаритные размеры / Масса" for r in out):
        out.append(["Габаритные размеры / Масса", *([""] * (width - 1))])
    for row in dims:
        if not row or not any(c.strip() for c in row):
            continue
        label = row[0].strip()
        val = next((c.strip() for c in row[1:] if c.strip()), "")
        padded = [label, val, *([""] * max(0, width - 2))]
        out.append(padded[:width])
    return out


@dataclass
class ManualDoc:
    stem: str
    title: str
    doc_title: str = ""
    torque: str = ""
    skus: list[str] = field(default_factory=list)
    warnings: str = ""
    contacts: str = ""
    summary: str = ""
    rotation_note: str = ""
    aux_table: list[list[str]] = field(default_factory=list)
    electrical_table: list[list[str]] = field(default_factory=list)
    function_table: list[list[str]] = field(default_factory=list)
    product_photo: str | None = None
    lead_photo: str | None = None
    diagram_photos: list[str] = field(default_factory=list)


# RU distributor contacts (replace PPTX China manufacturer block).
CONTACTS_HOOCON_RU = """\
Юридический адрес:
143440, Московская область, го. Красногорск, пгт. Путилково,
тер. Гринвуд, стр. 7, помещ. 98
Многоканальный телефон: 8 800 350-58-98
По вопросам сотрудничества: info@hoocon.ru
Отдел продаж: sales@hoocon.ru
www.hoocon.ru\
"""

CONTACTS_CHAMPION_BY = """\
Юридический адрес:
220030 Минск, пр-т Независимости 32А, пом.11
Многоканальный телефон: +375 29 372 6888
По вопросам сотрудничества: ichampiontech@yandex.ru\
"""

# Catalog media root (prepared wiring / dimensions per SKU family).
MEDIA_IMAGES = _repo_root() / "backend" / "media" / "product_images"

# Sheet-2 diagram template: same layout for all stems; assets/copy vary by model.
ROTATION_COPY_TERMINALS = (
    "<p>Заводская настройка: напряжение подаётся на клеммы 1 и 2, "
    "привод вращается против часовой стрелки.</p>"
    '<table class="data-table rotation-table"><tbody>'
    "<tr><td>Электропитание</td><td>Клеммы 1, 2</td><td>Клеммы 1, 3</td></tr>"
    "<tr><td>Направление вращения</td>"
    "<td>против часовой стрелки</td><td>по часовой стрелке</td></tr>"
    "</tbody></table>"
    '<p class="rotation-label">Поворотный переключатель</p>'
)
# Modulating DA2MU: rotary switch 0/1 vs rising control signal (not power terminals).
ROTATION_COPY_SIGNAL = (
    "<p>Установка направления вращения переключателем:</p>"
    '<table class="data-table rotation-table"><tbody>'
    "<tr><td>Положение переключателя</td><td>0</td><td>1</td></tr>"
    "<tr><td>При увеличении управляющего сигнала</td>"
    "<td>против часовой стрелки</td><td>по часовой стрелке</td></tr>"
    "</tbody></table>"
    '<p class="rotation-label">Поворотный переключатель</p>'
)
# DA4/6MU A/AS: same terminals table as ON/OFF plus motor-pin note from PDF.
ROTATION_COPY_TERMINALS_JUMPER = (
    "<p>Заводская настройка: напряжение подаётся на клеммы 1 и 2, "
    "привод вращается против часовой стрелки.</p>"
    '<table class="data-table rotation-table"><tbody>'
    "<tr><td>Электропитание</td><td>Клеммы 1, 2</td><td>Клеммы 1, 3</td></tr>"
    "<tr><td>Направление вращения</td>"
    "<td>против часовой стрелки</td><td>по часовой стрелке</td></tr>"
    "</tbody></table>"
    "<p>Направление вращения можно изменить с помощью перемычки на двигателе.</p>"
    '<p class="rotation-label">Поворотный переключатель</p>'
)
# EN MU/MQU D/DS: terminals + DIP S1 («Commutating switch»).
ROTATION_COPY_COMMUTATING = (
    "<p>Заводская настройка: напряжение подаётся на клеммы 1 и 3, "
    "привод вращается по часовой стрелке.</p>"
    '<table class="data-table rotation-table"><tbody>'
    "<tr><td>Электропитание</td><td>Клеммы 1, 2</td><td>Клеммы 1, 3</td></tr>"
    "<tr><td>Направление вращения</td>"
    "<td>против часовой стрелки</td><td>по часовой стрелке</td></tr>"
    "</tbody></table>"
    "<p>Направление вращения можно изменить DIP-переключателем S1.</p>"
)
ROTATION_LABEL_COMMUTATING = "Коммутирующий переключатель"
ROTATION_COPY_SCREW = (
    "<p>Направление вращения можно изменить винтом на верхней крышке.</p>"
)
ROTATION_COPY_ANGLE_LIMIT = (
    "<p>Механический упор угла поворота настраивается винтом "
    "на верхней крышке привода.</p>"
)
ROTATION_COPY_JUMPER = (
    "<p>Направление вращения можно изменить с помощью перемычки на двигателе.</p>"
)
ROTATION_COPY_SLIDER = (
    "<p>Направление вращения меняется ползунком переключателя "
    "на верхней крышке привода.</p>"
)
# SA FU/MU: reverse by mounting the actuator on the opposite side.
ROTATION_COPY_FLIP_SIDE = (
    "<p>Направление вращения задаётся ориентацией привода при монтаже "
    "(установка с противоположной стороны).</p>"
)

# RU caption under EN D/DS aux crop (Chinese callouts on the board drawing).
AUX_CAPTION_EN_MU_ONOFF = (
    "Винт M3 с крестообразным шлицем, полукруглая головка (a, b). "
    "Разъём двигателя."
)


@dataclass(frozen=True)
class DiagramProfile:
    """Per-manual sheet-2 diagram sources and rotation copy."""

    wiring_media: str | None  # basename under product_images, or None → curated shared
    dimensions_media: str  # basename under product_images
    wiring_overlays: bool  # RU overlays on curated ON/OFF wiring crop
    rotation_kind: str  # terminals | signal | terminals_jumper | screw | jumper | angle_limit | slider | dip | commutating
    rotation_image: bool  # show rotary-switch drawing
    # Optional local instruction PDF (under ``_инструкции-pdf/``) — overrides media crops.
    instruction_pdf: str | None = None
    # Sheet 1: aux table cols 1–3 + aux-diagram.png in cols 4–6 (else table spans 1–6).
    sheet1_aux_diagram: bool = False
    # Sheet 1: DIP / control-signal mode crop (A/AS modulating), below aux row.
    sheet1_dip_diagram: bool = False
    # Optional RU caption under aux diagram (e.g. translate CN callouts on EN crop).
    sheet1_aux_caption: str | None = None
    # Sheet 2 wiring: RU «Привод» / «Вспомогательный переключатель» above cropped PNG.
    wiring_ru_headers: bool = False


DIAGRAM_PROFILES: dict[str, DiagramProfile] = {
    # DA2MU ON/OFF — curated wiring crop + overlays; dims from catalog.
    "da2mu-d-ds": DiagramProfile(
        wiring_media=None,
        dimensions_media="da2mu230-ds-dimensions.webp",
        wiring_overlays=True,
        rotation_kind="terminals",
        rotation_image=True,
    ),
    "da2mu-a-as": DiagramProfile(
        wiring_media="da2mu24-as-wiring.webp",
        dimensions_media="da2mu24-as-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="signal",
        rotation_image=True,
    ),
    "da4-6mu-d-ds": DiagramProfile(
        wiring_media="da4mu230-ds-wiring.webp",
        dimensions_media="da4mu230-ds-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="terminals_jumper",
        rotation_image=True,
        instruction_pdf="da4_6mu-d_ds.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "da4-6mu-a-as": DiagramProfile(
        wiring_media="da4mu24-a-wiring.webp",
        dimensions_media="da4mu24-a-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="terminals_jumper",
        rotation_image=True,
        instruction_pdf="da4_6mu-a_as.pdf",
        sheet1_aux_diagram=True,
        sheet1_dip_diagram=True,
        wiring_ru_headers=True,
    ),
    "da3fu-d-ds": DiagramProfile(
        wiring_media="da3fu230-ds-wiring.webp",
        dimensions_media="da3fu230-ds-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="screw",
        rotation_image=False,
    ),
    "da5fu-d-ds": DiagramProfile(
        wiring_media="da5fu230-ds-wiring.webp",
        dimensions_media="da5fu230-ds-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="angle_limit",
        rotation_image=True,
        instruction_pdf="da5fu-d_ds.pdf",
    ),
    "da10-15-20fu24-230-d-ds": DiagramProfile(
        wiring_media="da10fu230-ds-wiring.webp",
        dimensions_media="da10fu230-ds-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="angle_limit",
        rotation_image=True,
        instruction_pdf="da10fu-d:ds.pdf",
        sheet1_aux_diagram=True,
    ),
    "da10-15-20fu24-a-as": DiagramProfile(
        wiring_media="da10fu24-a-wiring.webp",
        dimensions_media="da10fu24-a-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="slider",
        rotation_image=True,
        instruction_pdf="da10fu24-a:as.pdf",
        sheet1_aux_diagram=True,
    ),
    # EN→RU V24/V230 (separate voltage PDFs under ``_инструкции-pdf/EN/``).
    "da8-16-24-32mu24-a-as": DiagramProfile(
        wiring_media="da8mu24-a-wiring.webp",
        dimensions_media="da8mu24-a-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="dip",
        rotation_image=True,
        instruction_pdf="da8_16_24_32mu24-a_as.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "da8-16-24-32mu24-d-ds": DiagramProfile(
        wiring_media="da8mu24-d-wiring.webp",
        dimensions_media="da8mu24-d-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="commutating",
        rotation_image=True,
        instruction_pdf="da8_16_24_32mu24-d_ds.pdf",
        sheet1_aux_diagram=True,
        sheet1_aux_caption=AUX_CAPTION_EN_MU_ONOFF,
        wiring_ru_headers=True,
    ),
    "da8-16-24-32mu230-a-as": DiagramProfile(
        wiring_media="da8mu230-a-wiring.webp",
        dimensions_media="da8mu230-a-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="dip",
        rotation_image=True,
        instruction_pdf="da8_16_24_32mu230-a_as.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "da8-16-24-32mu230-d-ds": DiagramProfile(
        wiring_media="da8mu230-d-wiring.webp",
        dimensions_media="da8mu230-d-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="commutating",
        rotation_image=True,
        instruction_pdf="da8_16_24_32mu230-d_ds.pdf",
        sheet1_aux_diagram=True,
        sheet1_aux_caption=AUX_CAPTION_EN_MU_ONOFF,
        wiring_ru_headers=True,
    ),
    "da8-16-24mqu24-a-as": DiagramProfile(
        wiring_media="da8mqu24-a-wiring.webp",
        dimensions_media="da8mqu24-a-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="dip",
        rotation_image=True,
        instruction_pdf="da8_16_24mqu24-a_as.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "da8-16-24mqu230-a-as": DiagramProfile(
        wiring_media="da8mqu230-a-wiring.webp",
        dimensions_media="da8mqu230-a-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="dip",
        rotation_image=True,
        instruction_pdf="da8_16_24mqu230-a_as.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "da8-16-24mqu230-d-ds": DiagramProfile(
        wiring_media="da8mqu230-d-wiring.webp",
        dimensions_media="da8mqu230-d-dimensions.webp",
        wiring_overlays=False,
        rotation_kind="commutating",
        rotation_image=True,
        instruction_pdf="da8_16_24mqu230-d_ds.pdf",
        sheet1_aux_diagram=True,
        sheet1_aux_caption=AUX_CAPTION_EN_MU_ONOFF,
        wiring_ru_headers=True,
    ),
    # SA EN (combined 24+230, DS/DST) under ``_инструкции-pdf/EN/``.
    "sa3fu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa3fu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="thermal_saf72",
        rotation_image=True,
        instruction_pdf="sa3fu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "sa5fu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa5fu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="thermal_saf72",
        rotation_image=True,
        instruction_pdf="sa5fu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "sa10fu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa10fu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="thermal_saf72",
        rotation_image=True,
        instruction_pdf="sa10fu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "sa15fu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa15fu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="thermal_saf72",
        rotation_image=True,
        instruction_pdf="sa15fu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "sa7mu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa7mu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="flip_side",
        rotation_image=False,
        instruction_pdf="sa7mu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "sa10mu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa10mu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="flip_side",
        rotation_image=False,
        instruction_pdf="sa10mu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "sa15mu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa15mu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="flip_side",
        rotation_image=False,
        instruction_pdf="sa15mu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
    "sa30mu-ds-dst": DiagramProfile(
        wiring_media=None,
        dimensions_media="sa30mu-placeholder.webp",
        wiring_overlays=False,
        rotation_kind="flip_side",
        rotation_image=False,
        instruction_pdf="sa30mu-ds_dst.pdf",
        sheet1_aux_diagram=True,
        wiring_ru_headers=True,
    ),
}


def resolve_instruction_pdf(filename: str) -> Path:
    """Resolve an instruction PDF under ``_инструкции-pdf/{RU,EN}/`` or root."""
    root = _repo_root() / "_инструкции-pdf"
    name = Path(filename).name
    for sub in ("RU", "EN", ""):
        base = root / sub if sub else root
        candidate = base / name
        if candidate.is_file():
            return candidate
    raise FileNotFoundError(f"Instruction PDF missing: {root}/{{RU,EN}}/{name}")


def _find_media_file(basename: str) -> Path:
    """Resolve ``…_<basename>`` under product_images (prefer crop without title bar)."""
    needle = f"_{basename}"
    hits: list[Path] = []
    for path in MEDIA_IMAGES.rglob("*.webp"):
        if path.name.endswith(needle) or path.name == basename:
            hits.append(path)
    if not hits:
        raise FileNotFoundError(f"Catalog media not found: {basename}")
    if len(hits) == 1:
        return hits[0]
    # Prefer the file whose top rows are white (no baked-in black title).
    from PIL import Image
    import numpy as np

    scored: list[tuple[float, Path]] = []
    for path in hits:
        arr = np.asarray(Image.open(path).convert("RGB"))
        top = float(arr[: max(1, min(12, arr.shape[0]))].mean())
        scored.append((top, path))
    scored.sort(key=lambda item: item[0], reverse=True)
    return scored[0][1]


def _webp_to_png(src: Path, dest: Path, *, crop_title_banner: bool = False) -> None:
    from PIL import Image

    dest.parent.mkdir(parents=True, exist_ok=True)
    img = Image.open(src).convert("RGB")
    if crop_title_banner:
        img = _crop_dark_title_banner(img)
    img.save(dest, optimize=True)


def _crop_dark_title_banner(img: "Image.Image") -> "Image.Image":
    """Drop solid dark title bar at the top (duplicate of HTML h2.banner)."""
    import numpy as np

    arr = np.asarray(img.convert("RGB"))
    height = arr.shape[0]

    def row_mean(y: int) -> float:
        return float(arr[y].mean())

    y = 0
    while y < min(height, 8) and row_mean(y) > 200:
        y += 1
    start = y
    while y < min(height, 160) and row_mean(y) < 90:
        y += 1
    if y - start < 8:
        return img
    # Small pad under the bar; do not eat into the drawing.
    cut = min(height - 1, y + 2)
    return img.crop((0, cut, img.width, img.height))


def _crop_dark_bottom_banner(img: "Image.Image") -> "Image.Image":
    """Drop solid dark bar at the bottom (next-section banner / rule bleed)."""
    import numpy as np

    arr = np.asarray(img.convert("RGB"))
    height = arr.shape[0]

    def row_mean(y: int) -> float:
        return float(arr[y].mean())

    y = height - 1
    # Skip trailing white / light pad under the bar (trim_pad can be > 8 px).
    while y > max(0, height - 80) and row_mean(y) > 180:
        y -= 1
    end = y
    if end <= 0 or row_mean(end) >= 90:
        return img
    while y > max(0, height - 220) and row_mean(y) < 90:
        y -= 1
    # Thin hairline rules are often 2–6 px; keep a floor so we still strip them.
    if end - y < 2:
        return img
    # Small pad above the bar; do not eat into the drawing.
    cut = max(0, y - 2)
    return img.crop((0, 0, img.width, cut + 1))


def _crop_dark_top_banner(img: "Image.Image") -> "Image.Image":
    """Drop solid dark bar at the top (EN section title banner remnant)."""
    import numpy as np

    arr = np.asarray(img.convert("RGB"))
    height = arr.shape[0]

    def row_mean(y: int) -> float:
        return float(arr[y].mean())

    def row_dark_frac(y: int) -> float:
        return float((arr[y].mean(axis=1) < 90).mean())

    # Find the solid dark banner (skip leading white / anti-aliased edge).
    start: int | None = None
    for y in range(min(height - 1, 100)):
        if row_mean(y) < 100 and row_dark_frac(y) > 0.55:
            start = y
            break
    if start is None:
        return img
    y = start
    while y < min(height - 1, 220) and row_mean(y) < 100 and row_dark_frac(y) > 0.55:
        y += 1
    if y - start < 2:
        return img
    # Small pad below the bar; do not eat into top dimension ticks.
    cut = min(height - 1, y + 2)
    return img.crop((0, cut, img.width, height))


def _trim_white(img: "Image.Image", *, pad: int = 4) -> "Image.Image":
    import numpy as np

    arr = np.asarray(img.convert("RGB"))
    mask = arr.mean(axis=2) < 250
    ys, xs = np.where(mask)
    if len(xs) == 0:
        return img
    return img.crop(
        (
            max(0, int(xs.min()) - pad),
            max(0, int(ys.min()) - pad),
            min(img.width, int(xs.max()) + pad + 1),
            min(img.height, int(ys.max()) + pad + 1),
        ),
    )


def _pdf_extract_largest_image(
    pdf_path: Path,
    *,
    page_index: int,
    dest: Path,
    min_area: int = 40_000,
) -> bool:
    """Save the largest embedded raster from a PDF page (product hero photo).

    CMYK product shots often ship with a soft mask; extracting RGB without the
    mask yields ghosted/tiled sides (seen on ``da10fu24-a:as.pdf``).
    """
    import fitz
    from PIL import Image

    doc = fitz.open(pdf_path)
    try:
        page = doc[page_index]
        best_xref: int | None = None
        best_smask = 0
        best_area = 0
        for img in page.get_images(full=True):
            xref, smask, width, height = img[0], img[1], img[2], img[3]
            area = int(width) * int(height)
            if area < min_area or area <= best_area:
                continue
            best_xref = int(xref)
            best_smask = int(smask or 0)
            best_area = area
        if best_xref is None:
            return False
        pix = fitz.Pixmap(doc, best_xref)
        try:
            if pix.n - pix.alpha >= 4:  # CMYK etc.
                pix = fitz.Pixmap(fitz.csRGB, pix)
            if best_smask:
                pix = fitz.Pixmap(pix, fitz.Pixmap(doc, best_smask))
            if pix.alpha:
                rgba = Image.frombytes("RGBA", (pix.width, pix.height), pix.samples)
                img = Image.new("RGB", rgba.size, (255, 255, 255))
                img.paste(rgba, mask=rgba.split()[3])
            else:
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        finally:
            pix = None
        dest.parent.mkdir(parents=True, exist_ok=True)
        img.save(dest, optimize=True)
        return True
    finally:
        doc.close()


def _pdf_clip_png(
    pdf_path: Path,
    *,
    page_index: int,
    clip: tuple[float, float, float, float],
    dest: Path,
    scale: float = 4.0,
    crop_bottom_banner: bool = False,
    crop_top_banner: bool = False,
    trim_pad: int = 4,
    wipe_top_pt: float = 0.0,
    wipe_words: frozenset[str] | None = None,
    wipe_words_y: tuple[float, float] | None = None,
) -> None:
    import fitz
    from PIL import Image

    doc = fitz.open(pdf_path)
    try:
        page = doc[page_index]
        if wipe_words:
            y_lo, y_hi = wipe_words_y if wipe_words_y is not None else (0.0, 1e9)
            x_lo, _, x_hi, _ = clip
            for w in page.get_text("words"):
                text, x0, y0, x1, y1 = w[4], w[0], w[1], w[2], w[3]
                if text not in wipe_words:
                    continue
                if y0 < y_lo or y0 > y_hi or x1 < x_lo or x0 > x_hi:
                    continue
                page.add_redact_annot(
                    fitz.Rect(x0 - 1, y0 - 1, x1 + 1, y1 + 1),
                    fill=(1, 1, 1),
                )
            page.apply_redactions(images=0)
        pix = page.get_pixmap(
            matrix=fitz.Matrix(scale, scale),
            clip=fitz.Rect(*clip),
            alpha=False,
        )
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        if wipe_top_pt > 0:
            # Erase a solid top band (PDF points × scale → px).
            wipe_px = min(img.height, max(0, int(wipe_top_pt * scale)))
            if wipe_px:
                img.paste((255, 255, 255), (0, 0, img.width, wipe_px))
        if crop_top_banner:
            img = _crop_dark_top_banner(img)
        if crop_bottom_banner:
            img = _crop_dark_bottom_banner(img)
        img = _trim_white(img, pad=trim_pad)
        if crop_top_banner:
            img = _crop_dark_top_banner(img)
        if crop_bottom_banner:
            # Trim may re-include a thin rule; strip again after white trim.
            img = _crop_dark_bottom_banner(img)
        dest.parent.mkdir(parents=True, exist_ok=True)
        img.save(dest, optimize=True)
    finally:
        doc.close()


# Locked curated RU DIP table (Apple Creator Studio → white bg).
# Do not replace with PDF crops / screenshots without updating this digest.
DIP_DIAGRAM_LOCKED_NAME = "dip-diagram-control-signal-ru.png"
DIP_DIAGRAM_LOCKED_SHA256 = (
    "7a3180cd4ebf7a7b2f2f739e5a3e59f088b89b6b081cceedcc136d8712516f9e"
)

# Finished manuals — generators must not rewrite HTML/assets (unless --force).
LOCKED_MANUAL_STEMS: frozenset[str] = frozenset(
    {
        "da2mu-a-as",
        "da2mu-d-ds",
        "da3fu-d-ds",
        "da4-6mu-a-as",
        "da4-6mu-d-ds",
        "da5fu-d-ds",
        "da8-16-24-32mu24-a-as",
        "da8-16-24-32mu24-d-ds",
        "da8-16-24-32mu230-a-as",
        "da8-16-24-32mu230-d-ds",
        "da8-16-24mqu24-a-as",
        "da8-16-24mqu230-a-as",
        "da8-16-24mqu230-d-ds",
        "da10-15-20fu24-230-d-ds",
        "da10-15-20fu24-a-as",
        # SA — все FU/MU зафиксированы.
        "sa3fu-ds-dst",
        "sa5fu-ds-dst",
        "sa7mu-ds-dst",
        "sa10fu-ds-dst",
        "sa10mu-ds-dst",
        "sa15fu-ds-dst",
        "sa15mu-ds-dst",
        "sa30mu-ds-dst",
    }
)


def manual_stem_is_locked(stem: str) -> bool:
    return stem in LOCKED_MANUAL_STEMS


def _curated_assets_dir() -> Path:
    return _repo_root() / "docs" / "demo" / "manuals-ru" / "assets"


def _copy_locked_curated(
    *,
    stem_dir: Path,
    curated_name: str,
    dest_name: str,
    expected_sha256: str,
) -> None:
    """Copy a locked curated PNG into ``stem_dir``; refuse unexpected edits."""
    import hashlib

    curated = _curated_assets_dir() / curated_name
    if not curated.is_file():
        raise FileNotFoundError(f"Missing locked curated asset: {curated}")
    data = curated.read_bytes()
    digest = hashlib.sha256(data).hexdigest()
    if digest != expected_sha256:
        raise RuntimeError(
            f"Locked curated asset changed: {curated.name}\n"
            f"  expected sha256={expected_sha256}\n"
            f"  actual   sha256={digest}\n"
            "Update DIP_DIAGRAM_LOCKED_SHA256 only after intentional replace."
        )
    stem_dir.mkdir(parents=True, exist_ok=True)
    (stem_dir / dest_name).write_bytes(data)


def _copy_curated_dip_diagram(stem_dir: Path) -> None:
    """RU DIP table — locked curated (not EN PDF crop)."""
    _copy_locked_curated(
        stem_dir=stem_dir,
        curated_name=DIP_DIAGRAM_LOCKED_NAME,
        dest_name="dip-diagram.png",
        expected_sha256=DIP_DIAGRAM_LOCKED_SHA256,
    )


def _copy_curated_aux_diagram(stem_dir: Path) -> None:
    """RU aux / board layout (landscape, hub left) for EN MU/MQU manuals."""
    curated = _curated_assets_dir() / "aux-diagram-mu-board-ru.png"
    if not curated.is_file():
        raise FileNotFoundError(f"Missing curated aux diagram: {curated}")
    stem_dir.mkdir(parents=True, exist_ok=True)
    (stem_dir / "aux-diagram.png").write_bytes(curated.read_bytes())


def _materialize_da4_6mu_a_as_from_pdf(pdf_path: Path, stem_dir: Path) -> None:
    """Crop product + diagrams from ``_инструкции-pdf/da4_6mu-a_as.pdf``."""
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=(480, 130, 640, 300),
        dest=stem_dir / "product.png",
    )
    # Sheet 1: internal aux-switch layout beside factory-settings table.
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=(218, 68, 472, 168),
        dest=stem_dir / "aux-diagram.png",
        scale=5.0,
    )
    # DIP / control-signal mode — curated RU diagram (not EN PDF crop).
    _copy_curated_dip_diagram(stem_dir)
    # Lead: reuse product faceplate.
    product = stem_dir / "product.png"
    if product.is_file():
        (stem_dir / "lead.png").write_bytes(product.read_bytes())

    _pdf_clip_png(
        pdf_path,
        page_index=1,
        # Below EN «Actuator» / «Auxiliary switch»; RU headers in HTML.
        clip=(468, 98, 825, 198),
        dest=stem_dir / "wiring.png",
        crop_bottom_banner=True,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(468, 220, 825, 408),
        dest=stem_dir / "dimensions.png",
        crop_bottom_banner=True,
    )
    # Rotary-switch callout only (RU copy lives in HTML).
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(620, 430, 820, 550),
        dest=stem_dir / "rotation.png",
        scale=5.0,
    )


def _materialize_da4_6mu_d_ds_from_pdf(pdf_path: Path, stem_dir: Path) -> None:
    """Crop product + diagrams from ``_инструкции-pdf/da4_6mu-d_ds.pdf``."""
    # Page 0 product is a 3-tile composite; old clip x=505 cut the shaft clamp.
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=(478, 138, 640, 302),
        dest=stem_dir / "product.png",
        scale=5.0,
    )
    # Sheet 1: internal aux-switch / terminals diagram (beside factory-settings table).
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=(216, 80, 388, 175),
        dest=stem_dir / "aux-diagram.png",
        scale=5.0,
    )
    # Lead: same product family shot on page 1 (left).
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(46, 54, 146, 154),
        dest=stem_dir / "lead.png",
        scale=5.0,
    )

    _pdf_clip_png(
        pdf_path,
        page_index=1,
        # Below EN «Actuator» / «Auxiliary switch»; RU headers in HTML.
        clip=(468, 98, 825, 198),
        dest=stem_dir / "wiring.png",
        crop_bottom_banner=True,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        # Start below «Actuator Dimensions(mm)» title bar (HTML has its own banner).
        clip=(468, 232, 825, 410),
        dest=stem_dir / "dimensions.png",
        crop_bottom_banner=True,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(620, 430, 820, 555),
        dest=stem_dir / "rotation.png",
        scale=5.0,
    )


def _materialize_da5fu_d_ds_from_pdf(pdf_path: Path, stem_dir: Path) -> None:
    """Crop product + sheet-2 diagrams from portrait ``da5fu-d_ds.pdf``."""
    # Page 0 — product photo (left).
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=(45, 85, 175, 250),
        dest=stem_dir / "product.png",
        scale=5.0,
    )
    lead = stem_dir / "product.png"
    if lead.is_file():
        (stem_dir / "lead.png").write_bytes(lead.read_bytes())

    # Page 2 — wiring, dimensions, mechanical angle-limit (screw) drawing.
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 100, 390, 195),
        dest=stem_dir / "wiring.png",
        scale=4.0,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 238, 390, 385),
        dest=stem_dir / "dimensions.png",
        scale=4.0,
        crop_bottom_banner=True,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 405, 390, 545),
        dest=stem_dir / "rotation.png",
        scale=4.0,
    )


def _materialize_da10_15_20fu24_230_d_ds_from_pdf(pdf_path: Path, stem_dir: Path) -> None:
    """Crop diagrams from portrait ``da10fu-d:ds.pdf`` (shared 10/15/20 D/DS)."""
    # Full embedded product photo (page 0) — red-mask crop cut silver sides/shadow.
    if not _pdf_extract_largest_image(
        pdf_path, page_index=0, dest=stem_dir / "product.png"
    ):
        _pdf_clip_png(
            pdf_path,
            page_index=0,
            clip=(55, 75, 160, 275),
            dest=stem_dir / "product.png",
            scale=5.0,
        )
    product = stem_dir / "product.png"
    if product.is_file():
        (stem_dir / "lead.png").write_bytes(product.read_bytes())

    # Page 2 — wiring (actuator + aux), dimensions, mechanical angle-limit.
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 95, 390, 200),
        dest=stem_dir / "wiring.png",
        scale=4.0,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 225, 390, 385),
        dest=stem_dir / "dimensions.png",
        scale=4.0,
        crop_bottom_banner=True,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 405, 390, 545),
        dest=stem_dir / "rotation.png",
        scale=4.0,
    )
    # Page 3 — aux-switch schematic beside factory-settings table.
    _pdf_clip_png(
        pdf_path,
        page_index=3,
        clip=(248, 72, 390, 195),
        dest=stem_dir / "aux-diagram.png",
        scale=5.0,
    )


def _materialize_da10_15_20fu24_a_as_from_pdf(pdf_path: Path, stem_dir: Path) -> None:
    """Crop sheet-1/2 diagrams from portrait ``da10fu24-a:as.pdf`` (shared 10/15/20)."""
    if not _pdf_extract_largest_image(
        pdf_path, page_index=0, dest=stem_dir / "product.png"
    ):
        _pdf_clip_png(
            pdf_path,
            page_index=0,
            clip=(45, 70, 180, 295),
            dest=stem_dir / "product.png",
            scale=5.0,
        )
    product = stem_dir / "product.png"
    if product.is_file():
        (stem_dir / "lead.png").write_bytes(product.read_bytes())

    # Page 2 — wiring, dimensions, rotation-direction slider on cover.
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 90, 390, 200),
        dest=stem_dir / "wiring.png",
        scale=4.0,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 228, 390, 375),
        dest=stem_dir / "dimensions.png",
        scale=4.0,
        crop_bottom_banner=True,
    )
    _pdf_clip_png(
        pdf_path,
        page_index=2,
        clip=(40, 400, 390, 530),
        dest=stem_dir / "rotation.png",
        scale=4.0,
    )
    # Page 3 — aux-switch schematic beside factory-settings table.
    # Slight left inset removes grey table-edge bleed; keep S1 / YEL.
    _pdf_clip_png(
        pdf_path,
        page_index=3,
        clip=(248, 72, 390, 195),
        dest=stem_dir / "aux-diagram.png",
        scale=5.0,
    )


def _materialize_landscape_en_mu_from_pdf(
    pdf_path: Path,
    stem_dir: Path,
    *,
    modulating: bool,
) -> None:
    """Crop product + diagrams from landscape A4×2 EN MU/MQU manuals.

    Layout matches DA4/6 RU PDFs (sheet1 aux+product; sheet2 wiring/dims/rotation).
    """
    # Hero: full body + hub; right edge before black SKU/torque plate.
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=(482, 118, 662, 302),
        dest=stem_dir / "product.png",
        scale=5.0,
        trim_pad=20,
    )

    # Sheet-2 lead: page-1 product (no SKU plate / no section banner).
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(42, 50, 150, 148),
        dest=stem_dir / "lead.png",
        scale=5.0,
    )
    if not (stem_dir / "lead.png").is_file():
        product = stem_dir / "product.png"
        if product.is_file():
            (stem_dir / "lead.png").write_bytes(product.read_bytes())

    if modulating:
        # A/AS: curated RU board (DIP + terminals 1–5 / 21–26).
        _copy_curated_aux_diagram(stem_dir)
        # Sheet 2 «Dial switch set» — curated RU DIP table.
        _copy_curated_dip_diagram(stem_dir)
    else:
        # D/DS: ON/OFF board from PDF (terminals 1–3 / 21–26, no DIP bank).
        _pdf_clip_png(
            pdf_path,
            page_index=0,
            clip=(226, 56, 372, 180),
            dest=stem_dir / "aux-diagram.png",
            scale=5.0,
            trim_pad=12,
        )
        # Sheet 2: board + S1 (EN «Commutating switch» → RU label under img).
        _pdf_clip_png(
            pdf_path,
            page_index=1,
            clip=(618, 432, 825, 545),
            dest=stem_dir / "rotation.png",
            scale=5.0,
            trim_pad=16,
        )

    # Wiring: below EN «Actuator» / «Auxiliary switch» (RU headers in HTML).
    # Start above diagram ink; redact EN titles so terminal/switch tops stay.
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(468, 88, 825, 198),
        dest=stem_dir / "wiring.png",
        crop_bottom_banner=True,
        trim_pad=12,
        wipe_words=frozenset({"Actuator", "Auxiliary", "switch"}),
        wipe_words_y=(80.0, 105.0),
    )
    # Dimensions: below EN «Actuator Dimensions(mm)» (RU banner in HTML).
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(468, 232, 825, 410),
        dest=stem_dir / "dimensions.png",
        crop_bottom_banner=True,
    )


def _materialize_sa_en_from_pdf(
    pdf_path: Path,
    stem_dir: Path,
    *,
    spring_return: bool,
) -> None:
    """Crop product + diagrams from landscape A4×2 EN SA FU/MU manuals.

    Combined 24 В + 230 В; sheet 2 has wiring + dims; SAFU also SAF72 block.
    """
    # Hero: leave left margin so body is not flush-clipped; FU keeps SAF72 accessory.
    if spring_return:
        product_clip = (450, 105, 690, 315)
        product_pad = 20
    else:
        # SAMU: extra left white so chassis edge is not flush in the hero slot.
        product_clip = (430, 95, 685, 320)
        product_pad = 48
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=product_clip,
        dest=stem_dir / "product.png",
        scale=5.0,
        trim_pad=product_pad,
    )

    # Lead: stop above page-1 «Technical specification» bar (~y 162).
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(42, 48, 158, 158),
        dest=stem_dir / "lead.png",
        scale=5.0,
        trim_pad=12,
        crop_bottom_banner=True,
    )
    if not (stem_dir / "lead.png").is_file():
        product = stem_dir / "product.png"
        if product.is_file():
            (stem_dir / "lead.png").write_bytes(product.read_bytes())

    # Sheet 1 aux schematic only (no factory-table bleed on the left).
    _pdf_clip_png(
        pdf_path,
        page_index=0,
        clip=(245, 58, 385, 178),
        dest=stem_dir / "aux-diagram.png",
        scale=5.0,
        trim_pad=10,
    )

    # Stop above «Actuator Dimensions(mm)» bar (SAF10/15 ~y 201–202; others ~y 208).
    wiring_y1 = 195.0 if spring_return else 198.0
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(468, 88, 825, wiring_y1),
        dest=stem_dir / "wiring.png",
        crop_bottom_banner=True,
        trim_pad=12,
        wipe_words=frozenset({"Actuator", "Auxiliary", "switch"}),
        wipe_words_y=(70.0, 100.0),
    )

    # Dims under «Actuator Dimensions(mm)».
    # Clip includes the EN title bar (stripped); stop above footer (~y 387).
    # Do not wipe_words here — white text holes break solid-banner detection.
    dims_y0 = 190.0
    dims_y1 = 382.0
    _pdf_clip_png(
        pdf_path,
        page_index=1,
        clip=(468, dims_y0, 825, dims_y1),
        dest=stem_dir / "dimensions.png",
        crop_top_banner=True,
        crop_bottom_banner=True,
        trim_pad=10,
    )

    if spring_return:
        # Drawings only — EN prose lives as RU copy in HTML.
        _pdf_clip_png(
            pdf_path,
            page_index=1,
            clip=(468, 445, 825, 590),
            dest=stem_dir / "thermal.png",
            scale=4.0,
            trim_pad=6,
        )


def ensure_diagram_assets(
    stem: str,
    out_dir: Path,
    *,
    force: bool = False,
) -> DiagramProfile:
    """Copy/convert catalog diagrams into ``assets/<stem>/`` for this manual."""
    profile = DIAGRAM_PROFILES.get(stem)
    if profile is None:
        raise KeyError(f"No DiagramProfile for stem {stem!r}")
    if manual_stem_is_locked(stem) and not force:
        # Keep published HTML/assets as-is.
        return profile
    assets_root = out_dir / "assets"
    stem_dir = assets_root / stem
    stem_dir.mkdir(parents=True, exist_ok=True)

    if profile.instruction_pdf:
        pdf_path = resolve_instruction_pdf(profile.instruction_pdf)
        if stem == "da4-6mu-d-ds":
            _materialize_da4_6mu_d_ds_from_pdf(pdf_path, stem_dir)
            return profile
        if stem == "da4-6mu-a-as":
            _materialize_da4_6mu_a_as_from_pdf(pdf_path, stem_dir)
            return profile
        if stem == "da5fu-d-ds":
            _materialize_da5fu_d_ds_from_pdf(pdf_path, stem_dir)
            return profile
        if stem == "da10-15-20fu24-a-as":
            _materialize_da10_15_20fu24_a_as_from_pdf(pdf_path, stem_dir)
            return profile
        if stem == "da10-15-20fu24-230-d-ds":
            _materialize_da10_15_20fu24_230_d_ds_from_pdf(pdf_path, stem_dir)
            return profile
        if stem in {
            "da8-16-24-32mu24-a-as",
            "da8-16-24-32mu230-a-as",
            "da8-16-24mqu24-a-as",
            "da8-16-24mqu230-a-as",
        }:
            _materialize_landscape_en_mu_from_pdf(
                pdf_path, stem_dir, modulating=True
            )
            return profile
        if stem in {
            "da8-16-24-32mu24-d-ds",
            "da8-16-24-32mu230-d-ds",
            "da8-16-24mqu230-d-ds",
        }:
            _materialize_landscape_en_mu_from_pdf(
                pdf_path, stem_dir, modulating=False
            )
            return profile
        if stem in {
            "sa3fu-ds-dst",
            "sa5fu-ds-dst",
            "sa10fu-ds-dst",
            "sa15fu-ds-dst",
        }:
            _materialize_sa_en_from_pdf(pdf_path, stem_dir, spring_return=True)
            return profile
        if stem in {
            "sa7mu-ds-dst",
            "sa10mu-ds-dst",
            "sa15mu-ds-dst",
            "sa30mu-ds-dst",
        }:
            _materialize_sa_en_from_pdf(pdf_path, stem_dir, spring_return=False)
            return profile

    # Dimensions — always from catalog; strip baked-in title bar if present.
    _webp_to_png(
        _find_media_file(profile.dimensions_media),
        stem_dir / "dimensions.png",
        crop_title_banner=True,
    )

    # Wiring — curated shared crop for DA2MU ON/OFF; else catalog.
    if profile.wiring_media is None:
        curated = _curated_assets_dir() / "wiring-schematic-onoff.png"
        if not curated.is_file():
            raise FileNotFoundError(f"Missing curated wiring: {curated}")
        (stem_dir / "wiring.png").write_bytes(curated.read_bytes())
    else:
        _webp_to_png(_find_media_file(profile.wiring_media), stem_dir / "wiring.png")

    # Rotary-switch drawing (shared curated asset when needed).
    if profile.rotation_image:
        curated_rot = _curated_assets_dir() / "rotation-switch-schematic-onoff.png"
        if not curated_rot.is_file():
            raise FileNotFoundError(f"Missing curated rotation: {curated_rot}")
        (stem_dir / "rotation.png").write_bytes(curated_rot.read_bytes())

    return profile


def _wiring_figure_html(
    stem: str,
    *,
    overlays: bool,
    ru_headers: bool = False,
) -> str:
    img = f'<img src="assets/{stem}/wiring.png" alt="Схема подключения">'
    headers = ""
    if ru_headers:
        headers = (
            '<div class="wiring-ru-headers" aria-hidden="true">'
            '<span class="wrh wrh-act">Привод</span>'
            '<span class="wrh wrh-aux">Вспомогательный переключатель</span>'
            "</div>"
        )
    board_cls = "wiring-board wiring-board-ru-headers" if ru_headers else "wiring-board"
    if not overlays:
        return (
            f'<figure class="diagram diagram-wide wiring-figure">'
            f'<div class="{board_cls}">{headers}{img}</div></figure>'
        )
    return (
        f'<figure class="diagram diagram-wide wiring-figure">'
        f'<div class="{board_cls}">'
        f"{headers}{img}"
        '<span class="wl wl-p2a">2-позиционное</span>'
        '<span class="wl wl-p3a">3-позиционное</span>'
        '<span class="wl wl-p2b">2-позиционное</span>'
        '<span class="wl wl-p3b">3-позиционное</span>'
        '<span class="wl wl-v24">AC/DC 24 В</span>'
        '<span class="wl wl-v230">AC 100…240 В</span>'
        '<span class="wl wl-rating">3(1.5) А AC 250 В</span>'
        "</div></figure>"
    )


def _rotation_figure_html(stem: str, profile: DiagramProfile) -> str:
    label_under_img = ""
    if profile.rotation_kind == "terminals":
        copy = ROTATION_COPY_TERMINALS
        banner = "Переключение направления вращения"
    elif profile.rotation_kind == "signal":
        copy = ROTATION_COPY_SIGNAL
        banner = "Переключение направления вращения"
    elif profile.rotation_kind == "terminals_jumper":
        copy = ROTATION_COPY_TERMINALS_JUMPER
        banner = "Переключение направления вращения"
    elif profile.rotation_kind == "commutating":
        copy = ROTATION_COPY_COMMUTATING
        banner = "Переключение направления вращения"
        label_under_img = ROTATION_LABEL_COMMUTATING
    elif profile.rotation_kind == "jumper":
        copy = ROTATION_COPY_JUMPER
        banner = "Переключение направления вращения"
    elif profile.rotation_kind == "angle_limit":
        copy = ROTATION_COPY_ANGLE_LIMIT
        banner = "Способ установки механического ограничения положения (угла поворота)"
    elif profile.rotation_kind == "slider":
        copy = ROTATION_COPY_SLIDER
        banner = "Изменение положения ручки переключателя"
    elif profile.rotation_kind == "dip":
        # EN MU/MQU A/AS: Dial switch set lives on sheet 2 (not sheet 1).
        banner = "Настройки DIP-переключателя"
        return (
            '<figure class="diagram diagram-wide diagram-rotation diagram-dip">'
            f'<h2 class="banner">{banner}</h2>'
            f'<img src="assets/{stem}/dip-diagram.png" alt="{banner}">'
            "</figure>"
        )
    elif profile.rotation_kind == "thermal_saf72":
        banner = "Термодатчик SAF72"
        copy = (
            "<p>Состоит из датчика окружающей среды (TS1) и канального "
            "датчика (TS2).</p>"
            "<p>TS1 размыкается при температуре окружающей среды выше "
            "72&nbsp;°C.</p>"
            "<p>TS2 размыкается при температуре в канале выше 72&nbsp;°C.</p>"
            "<p>Поставляется с исполнением «DST».</p>"
        )
        return (
            '<figure class="diagram diagram-wide diagram-rotation diagram-thermal">'
            f'<h2 class="banner">{banner}</h2>'
            '<div class="rotation-panel">'
            f'<div class="rotation-copy">{copy}</div>'
            f'<img src="assets/{stem}/thermal.png" alt="{banner}">'
            "</div></figure>"
        )
    elif profile.rotation_kind == "flip_side":
        copy = ROTATION_COPY_FLIP_SIDE
        banner = "Переключение направления вращения"
        return (
            '<figure class="diagram diagram-wide diagram-rotation">'
            f'<h2 class="banner">{banner}</h2>'
            f'<div class="rotation-copy">{copy}</div></figure>'
        )
    else:
        copy = ROTATION_COPY_SCREW
        banner = "Переключение направления вращения"
    if not profile.rotation_image:
        return (
            '<figure class="diagram diagram-wide diagram-rotation">'
            f'<h2 class="banner">{banner}</h2>'
            f'<div class="rotation-copy">{copy}</div></figure>'
        )
    img = f'<img src="assets/{stem}/rotation.png" alt="">'
    if label_under_img:
        media = (
            f'<div class="rotation-media">{img}'
            f'<p class="rotation-label">{html.escape(label_under_img)}</p></div>'
        )
    else:
        media = img
    return (
        '<figure class="diagram diagram-wide diagram-rotation">'
        f'<h2 class="banner">{banner}</h2>'
        f'<div class="rotation-panel"><div class="rotation-copy">{copy}</div>'
        f"{media}</div></figure>"
    )


def curated_diagrams_html(stem: str, profile: DiagramProfile) -> str:
    """Locked sheet-2 diagrams block (wiring → dimensions → rotation)."""
    parts = [
        _wiring_figure_html(
            stem,
            overlays=profile.wiring_overlays,
            ru_headers=profile.wiring_ru_headers,
        ),
        (
            '<figure class="diagram diagram-wide">'
            '<h2 class="banner">Габаритные размеры привода (мм)</h2>'
            f'<img src="assets/{stem}/dimensions.png" alt="">'
            "</figure>"
        ),
        _rotation_figure_html(stem, profile),
    ]
    return "".join(parts)


def _save_pic(sh, assets: Path, stem: str, idx: int) -> str | None:
    try:
        blob = sh.image.blob
        ext = sh.image.ext or "png"
    except Exception:
        return None
    w, h = int(sh.width or 0), int(sh.height or 0)
    if w < 50000 or h < 50000:
        return None
    area = w * h
    # Skip dust/glyph fragments; keep product & diagram images.
    if area < 900_000_000 and len(blob) < 2500:
        return None
    name = f"{stem}_{idx:03d}.{ext}"
    (assets / name).write_bytes(blob)
    return name


def extract_manual(
    prs: Presentation,
    assets: Path,
    stem: str,
    page_title: str,
) -> ManualDoc:
    doc = ManualDoc(stem=stem, title=page_title)
    img_i = 0
    slide_pics: list[list[tuple[int, str]]] = [[], []]
    dims_tables: list[list[list[str]]] = []

    for si, slide in enumerate(prs.slides):
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        for sh, _l, _t in iter_shapes(slide.shapes):
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_i += 1
                name = _save_pic(sh, assets, stem, img_i)
                if name and si < 2:
                    area = int(sh.width or 0) * int(sh.height or 0)
                    slide_pics[si].append((area, name))
            if sh.has_text_frame:
                raw = _shape_text(sh)
                if raw:
                    texts.append(translate(raw))
            if sh.has_table:
                matrix = _table_matrix(sh)
                if any(any(c for c in row) for row in matrix):
                    tables.append(matrix)

        for t in texts:
            low = t.lower()
            if "запрещается" in low or t.strip().startswith("1."):
                if len(t) > len(doc.warnings):
                    doc.warnings = t
            elif "ningbo" in low or "hoocon@" in low or "тел." in low:
                doc.contacts = t
            elif "направление вращения" in low and len(t) < 220:
                doc.rotation_note = t
            elif ("руководство" in low or "привод" in low) and len(t) < 140:
                if "руководство" in low or not doc.doc_title:
                    doc.doc_title = re.sub(r"\s*\n\s*", " — ", t).strip(" —")
            elif re.search(r"\d+\s*нм", t, re.I) and (
                "управл" in low
                or "on/off" in low
                or "пропорционал" in low
                or "момент" in low
                or "время" in low
            ):
                doc.summary = t
            elif re.fullmatch(r"\d+(?:\s*/\s*\d+)*\s*нм", t.strip(), re.I):
                doc.torque = t.strip()

        blob = "\n".join(texts)
        for matrix in tables:
            for row in matrix:
                blob += "\n" + " ".join(row)
        skus: list[str] = []
        for s in _SKU_TOKEN.findall(blob):
            u = s.upper()
            if u not in skus and len(u) >= 5:
                skus.append(u)
        if skus:
            doc.skus = skus[:12]

        for matrix in tables:
            if _is_banner_row(matrix):
                continue
            joined = " ".join(c for row in matrix for c in row).lower()
            n_rows = len(matrix)
            n_cols = len(matrix[0]) if matrix else 0
            if n_rows >= 2 and n_cols >= 2 and (
                "переключатель" in joined or "замкнуто" in joined or "клемм" in joined
            ):
                if not doc.aux_table or n_rows >= len(doc.aux_table):
                    doc.aux_table = matrix
            elif n_rows >= 10:
                doc.function_table = matrix
            elif _is_dimensions_weight_table(matrix):
                dims_tables.append(matrix)
            elif n_rows >= 2 and n_cols >= 2 and (
                "напряжение" in joined or "потребляем" in joined or "сечение" in joined
            ):
                doc.electrical_table = matrix
            elif si == 1 and 2 <= n_rows <= 8 and n_cols >= 2 and not doc.electrical_table:
                doc.electrical_table = matrix

    for dims in dims_tables:
        doc.function_table = _merge_dimensions_into_function(doc.function_table, dims)
    if slide_pics[0]:
        slide_pics[0].sort(reverse=True)
        doc.product_photo = slide_pics[0][0][1]
        # Sheet 2 summary uses the same hero product shot as sheet 1.
        doc.lead_photo = doc.product_photo
    if slide_pics[1]:
        slide_pics[1].sort(reverse=True)
        pics = [n for _a, n in slide_pics[1][:4]]
        if len(pics) >= 2:
            # Skip first large image (often a duplicate/photo fragment).
            mid = pics[:-1]
            doc.diagram_photos = mid[1:] if len(mid) > 1 else mid
        else:
            doc.diagram_photos = pics
        if not doc.lead_photo:
            doc.lead_photo = pics[0] if pics else None

    if not doc.doc_title:
        doc.doc_title = "Руководство по эксплуатации"
    if not doc.torque and doc.summary:
        m = re.search(r"(\d+(?:\s*/\s*\d+)*)\s*нм", doc.summary, re.I)
        if m:
            doc.torque = m.group(0)
    _normalize_aux_edition_copy(doc)
    _normalize_multi_nm_running_time(doc)
    _normalize_da5fu_d_ds(doc)
    _normalize_da10_15_20fu24_a_as(doc)
    return doc


def _normalize_multi_nm_running_time(doc: ManualDoc) -> None:
    """Multi-Nm FU lead bullets: distinct running times from PDF."""
    if not doc.summary:
        return
    if doc.stem == "da10-15-20fu24-a-as":
        # 10 Нм → <110 с; 15/20 Нм → <150 с.
        doc.summary = doc.summary.replace(
            "Время поворота: < 110 с; время возврата пружины: < 25 с",
            "Время поворота: < 110 с / 150 с; время возврата пружины: < 25 с",
        )
    elif doc.stem == "da10-15-20fu24-230-d-ds":
        # 10 Нм → <100 с; 15/20 Нм → <150 с (PDF; not 170).
        doc.summary = doc.summary.replace(
            "Время поворота: < 100 с / 150 с / 170 с; время возврата пружины: < 25 с",
            "Время поворота: < 100 с / 150 с; время возврата пружины: < 25 с",
        )


def _normalize_aux_edition_copy(doc: ManualDoc) -> None:
    """D/DS manuals must not keep AS boilerplate from shared PPTX copy (and vice versa)."""
    if doc.stem.endswith("-d-ds"):
        repl = ("«AS»", "«DS»")
        banner_from, banner_to = (
            "Настройка вспомогательных переключателей (модели AS)",
            "Настройка вспомогательных переключателей (модели DS)",
        )
    elif doc.stem.endswith("-a-as"):
        # Keep manufacturer «DS» wording on DA2MU A/AS (1 group); only fix AS/DS letter
        # when the sentence is clearly the 2-group AS boilerplate mistyped as DS.
        repl = None
        banner_from, banner_to = (
            "Настройка вспомогательных переключателей (модели DS)",
            "Настройка вспомогательных переключателей (модели AS)",
        )
    else:
        return

    def fix(text: str) -> str:
        if not text:
            return text
        out = text.replace(banner_from, banner_to)
        if repl:
            out = out.replace(repl[0], repl[1])
        return out

    doc.summary = fix(doc.summary)
    doc.warnings = fix(doc.warnings)
    doc.rotation_note = fix(doc.rotation_note)
    _normalize_onoff_mu_summary(doc)


def _normalize_onoff_mu_summary(doc: ManualDoc) -> None:
    """DA4/6MU D/DS PPTX often inherits A/AS Y/U lines — PDF is ON/OFF 2/3-point."""
    if doc.stem != "da4-6mu-d-ds" or not doc.summary:
        return
    # Drop modulating signal bullets (wrong edition).
    doc.summary = re.sub(
        r"[ \t]*Упр\.\s*сигнал\s*Y\s*:.*?(?=\n|$)",
        "",
        doc.summary,
        flags=re.I,
    )
    doc.summary = re.sub(
        r"[ \t]*Обратная связь\s*U\s*:.*?(?=\n|$)",
        "",
        doc.summary,
        flags=re.I,
    )
    control = "Управление: вкл./выкл. (ON/OFF), 2- или 3-позиционное"
    if not re.search(r"(?m)^Управление\s*:", doc.summary):
        if re.search(r"Номинальное напряжение\s*:", doc.summary, re.I):
            doc.summary = re.sub(
                r"(Номинальное напряжение\s*:[^\n]+)",
                rf"\1\n{control}",
                doc.summary,
                count=1,
                flags=re.I,
            )
        else:
            doc.summary = f"{doc.summary.rstrip()}\n{control}"
    doc.summary = re.sub(r"[ \t]+\n", "\n", doc.summary)
    doc.summary = re.sub(r"\n{3,}", "\n\n", doc.summary)


_POWER_LONG_RE = re.compile(
    r"^(\d+(?:,\d+)?)\s*Вт\s+при\s+номинальном(?:\s+крутящем)?\s+моменте\s*/\s*"
    r"(\d+(?:,\d+)?)\s*Вт\s+(?:в\s+режиме\s+)?удержания$",
    re.I,
)


def _compact_power_text(text: str) -> str:
    """Short PDF-style power cell (two lines) — full phrase blows up narrow cols."""
    m = _POWER_LONG_RE.match(text.strip())
    if not m:
        return text
    return f"{m.group(1)} Вт под нагрузкой\n{m.group(2)} Вт удержание"


def _format_tech_label(text: str) -> str:
    """Break long first-column labels onto intentional lines."""
    text = re.sub(
        r"^(Габаритные размеры)\s*\((Д\s*×\s*Ш\s*×\s*В)\)$",
        r"\1\n(\2)",
        text.strip(),
        flags=re.I,
    )
    return text


def _escape_cell(text: str) -> str:
    return html.escape(_format_tech_label(_compact_power_text(text))).replace("\n", "<br>")


def _html_table(matrix: list[list[str]], *, stretch_rows: bool = False) -> str:
    if not matrix:
        return ""
    rows_html: list[str] = []
    for i, row in enumerate(matrix):
        if not any(c.strip() for c in row):
            continue
        n = len(row)
        label = row[0].strip()
        rest = row[1:]
        rest_filled = [c for c in rest if c.strip()]
        section = bool(label) and not rest_filled
        cells: list[str] = []
        if section:
            cells.append(
                f'<td class="row-section" colspan="{n}">{_escape_cell(label)}</td>'
            )
        elif len(rest) > 1 and len({c.strip() for c in rest}) == 1:
            # Same value in every data column — one colspan (clean borders).
            cells.append(f"<td>{_escape_cell(label)}</td>")
            cells.append(
                f'<td colspan="{len(rest)}">{_escape_cell(rest[0].strip())}</td>'
            )
        elif len(rest_filled) == 1 and len(rest) > 1:
            # One shared value for all data columns — merge empties.
            cells.append(f"<td>{_escape_cell(label)}</td>")
            cells.append(
                f'<td colspan="{len(rest)}">{_escape_cell(rest_filled[0])}</td>'
            )
        else:
            cells.append(f"<td>{_escape_cell(label)}</td>")
            i = 0
            while i < len(rest):
                if not rest[i].strip():
                    i += 1
                    continue
                val = rest[i].strip()
                j = i + 1
                # Merge consecutive identical values (e.g. 24×3 | 230×3).
                while j < len(rest) and rest[j].strip() == val:
                    j += 1
                span = j - i
                attr = f' colspan="{span}"' if span > 1 else ""
                cells.append(f"<td{attr}>{_escape_cell(val)}</td>")
                i = j
        rows_html.append("<tr>" + "".join(cells) + "</tr>")
    if not rows_html:
        return ""
    attrs = ' class="data-table"'
    if stretch_rows:
        n_rows = len(rows_html)
        attrs += f' style="--tech-rows:{n_rows}"'
        # Short electrical tables must not shrink — wrapped power cells would clip
        # the following «Сечение…» row under overflow:hidden.
        shrink = 0 if n_rows <= 6 else 1
        table = f"<table{attrs}><tbody>{''.join(rows_html)}</tbody></table>"
        return (
            f'<div class="tech-table-slot" style="flex:{n_rows} {shrink} auto">'
            f"{table}</div>"
        )
    return f"<table{attrs}><tbody>{''.join(rows_html)}</tbody></table>"


def _sku_voltage_groups(skus: list[str]) -> tuple[list[str], list[str]]:
    """Split SKUs into AC/DC 24 V vs AC 100…240 V columns."""
    v230 = [s for s in skus if "230" in s.upper()]
    v24 = [s for s in skus if s not in v230]
    return v24, v230


PROTECTION_CLASS_III = "III (безопасное сверхнизкое напряжение)"
PROTECTION_CLASS_II = "II (все изолировано / полная изоляция)"


@dataclass(frozen=True)
class VoltageTemplate:
    """Single-voltage manual shell (24 V XOR 230 V) for EN→RU translations.

    Manufacturer PDFs for MU8… / MQU often ship as separate ``*mu24*`` /
    ``*mu230*`` files — use these templates, not a combined 24|230 ТТХ.
    Shared A4×2 layout stays in ``HTML_SHELL`` / ``TEMPLATES.md``.
    """

    id: str  # v24 | v230
    label_ru: str
    sku_token: str  # infix in SKU: 24 | 230
    nominal_voltage: str
    protection_class: str
    lead_voltage: str
    example_skus: tuple[str, ...]
    html_stem: str
    page_title: str
    notes_ru: str


VOLTAGE_TEMPLATES: dict[str, VoltageTemplate] = {
    "v24": VoltageTemplate(
        id="v24",
        label_ru="AC/DC 24 В",
        sku_token="24",
        nominal_voltage="AC/DC 24 В, 50/60 Гц",
        protection_class=PROTECTION_CLASS_III,
        lead_voltage="Номинальное напряжение: AC/DC 24 В",
        example_skus=("DA…MU24-A/AS", "DA…MU24-D/DS"),
        html_stem="template-v24",
        page_title="Шаблон V24 — AC/DC 24 В (RU)",
        notes_ru=(
            "Одно напряжение на весь PDF. Класс защиты III (SELV). "
            "Колонки ТТХ — по Nm / исполнениям внутри 24 В, без колонки 230 В. "
            "Схемы подключения — из соответствующего *24* PDF."
        ),
    ),
    "v230": VoltageTemplate(
        id="v230",
        label_ru="AC 100…240 В",
        sku_token="230",
        nominal_voltage="AC 100…240 В, 50/60 Гц",
        protection_class=PROTECTION_CLASS_II,
        lead_voltage="Номинальное напряжение: AC 100…240 В",
        example_skus=("DA…MU230-A/AS", "DA…MU230-D/DS"),
        html_stem="template-v230",
        page_title="Шаблон V230 — AC 100…240 В (RU)",
        notes_ru=(
            "Одно напряжение на весь PDF. Класс защиты II (полная изоляция). "
            "Колонки ТТХ — по Nm / исполнениям внутри 230 В, без колонки 24 В. "
            "Схемы подключения — из соответствующего *230* PDF (не копировать 24 В)."
        ),
    ),
}


def voltage_template_for_sku(sku: str) -> VoltageTemplate:
    """Map one SKU code to the V24 or V230 shell."""
    return VOLTAGE_TEMPLATES["v230" if _sku_is_230v(sku) else "v24"]


def voltage_template_for_skus(skus: list[str]) -> VoltageTemplate | None:
    """Return V24/V230 when all SKUs share one voltage; else None (mixed manual)."""
    if not skus:
        return None
    ids = {voltage_template_for_sku(s).id for s in skus}
    if len(ids) != 1:
        return None
    return VOLTAGE_TEMPLATES[next(iter(ids))]


# DA10/15/20 FU power from published manuals (A/AS ≠ D/DS for 10 Нм).
# Compact two-line form (PDF: «под нагрузкой» / «удержание») — fits narrow ТТХ cols.
_POWER_FU24_A_AS_BY_NM: dict[int, str] = {
    10: "6 Вт под нагрузкой\n1,5 Вт удержание",
    15: "10 Вт под нагрузкой\n3 Вт удержание",
    20: "10 Вт под нагрузкой\n3 Вт удержание",
}
_POWER_FU_D_DS_BY_NM: dict[int, str] = {
    10: "5 Вт под нагрузкой\n3 Вт удержание",
    15: "10 Вт под нагрузкой\n3 Вт удержание",
    20: "10 Вт под нагрузкой\n3 Вт удержание",
}
# Motor / spring return (PDF D/DS); 15 and 20 share 150 с (not 170).
_RUNNING_FU_D_DS_BY_NM: dict[int, str] = {
    10: "< 100 с / < 25 с",
    15: "< 150 с / < 25 с",
    20: "< 150 с / < 25 с",
}
_MASS_FU_D_DS_BY_NM: dict[int, str] = {
    10: "< 2,3 кг",
    15: "< 2,5 кг",
    20: "< 2,5 кг",
}
# DA10/15/20 FU24 A/AS — from da10/15/20fu24-a:as.pdf (15/20 share 150 с).
_RUNNING_FU24_A_AS_BY_NM: dict[int, str] = {
    10: "< 110 с / < 25 с",
    15: "< 150 с / < 25 с",
    20: "< 150 с / < 25 с",
}
_AREA_FU24_A_AS_BY_NM: dict[int, str] = {
    10: "< 1 м²",
    15: "< 1,5 м²",
    20: "< 2 м²",
}
_MASS_FU24_A_AS_BY_NM: dict[int, str] = {
    10: "< 2,6 кг",
    15: "< 2,6 кг",
    20: "< 2,6 кг",
}


def _sku_is_230v(sku: str) -> bool:
    return "230" in sku.upper()


def _protection_class_for_sku(sku: str) -> str:
    """24 V → class III (SELV); 230 V → class II (totally insulated)."""
    return voltage_template_for_sku(sku).protection_class


def build_voltage_template_doc(vt: VoltageTemplate) -> ManualDoc:
    """Placeholder ManualDoc for the V24/V230 HTML shells (EN→RU starter)."""
    sku = vt.example_skus[0].replace("…", "N")
    edition_letter = "AS" if sku.upper().endswith("A/AS") or "-A" in sku.upper() else "DS"
    summary = "\n".join(
        [
            "Пропорциональное / вкл.-выкл. управление — подставить из PDF",
            "Для управления воздушными регулирующими заслонками в системе ОВК",
            "Крутящий момент: … Нм",
            "Время поворота: < … с; время возврата пружины: < … с",
            vt.lead_voltage,
            "Управление: … (из PDF / глоссарий Belimo RU)",
            f"Исполнение «{edition_letter}»: … (группы вспом. переключателей)",
        ]
    )
    warnings = (
        "1. Привод заслонки не разрешается использовать вне указанной области "
        "применения, особенно на воздушных судах.\n"
        "2. Корпус привода может вскрываться только изготовителем.\n"
        "3. Устройство содержит электрические и электронные компоненты и не "
        "подлежит утилизации с бытовыми отходами.\n"
        "4. Необходимо соблюдать все местные действующие нормы и требования."
    )
    elec = [
        ["Номинальное напряжение", vt.nominal_voltage],
        ["Потребляемая мощность", "… Вт под нагрузкой\n… Вт удержание"],
        ["Сечение подключаемых проводов", "0,5 мм²"],
    ]
    func = [
        ["Функциональные параметры", ""],
        ["Площадь обслуживаемой заслонки", "< … м²"],
        ["Направление вращения", "для монтажа с противоположной стороны"],
        ["Ручное управление", "… (из PDF)"],
        ["Угол поворота", "Макс. 95°"],
        ["Уровень звуковой мощности", "макс. … дБ(А) …"],
        ["Индикация положения", "Механический указатель"],
        ["Условия эксплуатации", ""],
        ["Класс защиты", vt.protection_class],
        ["Степень защиты корпуса", "IP…"],
        ["Температура окружающей среды", "–20…+50 °C"],
        ["Температура хранения", "–40…+70 °C"],
        ["Испытание на влажность", "5…95 % относительной влажности"],
        ["Габаритные размеры / Масса", ""],
        ["Габаритные размеры (Д × Ш × В)", "См. «Габаритные размеры»"],
        ["Длина вала заслонки", "> … мм"],
        ["Диаметр вала заслонки", "круглый … мм, квадратный … мм"],
        ["Масса", "< … кг"],
    ]
    return ManualDoc(
        stem=vt.html_stem,
        title=vt.page_title,
        doc_title="Руководство по эксплуатации — шаблон",
        torque="… Нм",
        skus=[sku],
        warnings=warnings,
        summary=summary,
        electrical_table=elec,
        function_table=func,
        aux_table=[
            ["Переключатель a", "Клеммы S1, S2", "Клеммы S1, S3"],
            ["…°", "Замкнуто", "Разомкнуто"],
            ["…°", "Разомкнуто", "Замкнуто"],
            ["Переключатель b", "Клеммы S4, S5", "Клеммы S4, S6"],
            ["…°", "Разомкнуто", "Замкнуто"],
            ["…°", "Замкнуто", "Разомкнуто"],
        ],
    )


def write_voltage_template_shells(out_dir: Path) -> list[Path]:
    """Write ``template-v24.html`` / ``template-v230.html`` reference shells."""
    written: list[Path] = []
    for vt in VOLTAGE_TEMPLATES.values():
        doc = build_voltage_template_doc(vt)
        body = render_grid(doc, diagram_profile=None)
        note = (
            f'<p class="template-voltage-note"><strong>{html.escape(vt.label_ru)}</strong>'
            f" — {html.escape(vt.notes_ru)} "
            f"Канон: <code>TEMPLATES.md</code> § V24/V230; "
            f"термины: <code>docs/tech-copy-belimo-ru.md</code>.</p>"
        )
        body = body.replace(
            '<section class="sheet sheet-page1">',
            f'{note}<section class="sheet sheet-page1">',
            1,
        )
        path = out_dir / f"{vt.html_stem}.html"
        path.write_text(
            HTML_SHELL.format(title=html.escape(vt.page_title), body=body),
            encoding="utf-8",
        )
        written.append(path)
    return written


def _sku_torque_nm(sku: str) -> int | None:
    m = re.search(r"DA(\d+)(?:FU|MU)", sku, re.I)
    return int(m.group(1)) if m else None


def _power_consumption_for_sku(sku: str) -> str | None:
    """Per-SKU power for DA10/15/20 FU (PPTX often copies one value across Nm)."""
    nm = _sku_torque_nm(sku)
    if nm is None or "FU" not in sku.upper():
        return None
    u = sku.upper()
    if re.search(r"-A(?:/AS)?$|-AS$", u):
        return _POWER_FU24_A_AS_BY_NM.get(nm)
    if re.search(r"-D(?:/DS)?$|-DS$", u):
        return _POWER_FU_D_DS_BY_NM.get(nm)
    return None


def _normalize_protection_class_rows(
    matrix: list[list[str]],
    skus: list[str],
    *,
    per_sku: bool,
) -> list[list[str]]:
    """Rewrite «Класс защиты» from SKU voltage (PPTX often mislabels 24 V as II)."""
    if not matrix or not skus:
        return matrix
    out: list[list[str]] = []
    for row in matrix:
        if not row:
            out.append(row)
            continue
        label = row[0].strip()
        if not label.startswith("Класс защиты"):
            out.append(row)
            continue
        rest = row[1:] if len(row) > 1 else []
        n = len(rest)
        if per_sku and n == len(skus):
            out.append(["Класс защиты"] + [_protection_class_for_sku(s) for s in skus])
            continue
        if n == 2:
            # Voltage layout: col1 = 24 V, col2 = 230 V.
            out.append(["Класс защиты", PROTECTION_CLASS_III, PROTECTION_CLASS_II])
            continue
        if n == 1:
            v24, v230 = _sku_voltage_groups(skus)
            if v24 and not v230:
                val = PROTECTION_CLASS_III
            elif v230 and not v24:
                val = PROTECTION_CLASS_II
            else:
                val = rest[0].strip() or PROTECTION_CLASS_III
            out.append(["Класс защиты", val])
            continue
        if n == len(skus):
            out.append(["Класс защиты"] + [_protection_class_for_sku(s) for s in skus])
            continue
        out.append(["Класс защиты"] + rest)
    return out


def _normalize_power_consumption_rows(
    matrix: list[list[str]],
    skus: list[str],
) -> list[list[str]]:
    """Split «Потребляемая мощность» per SKU when DA10/15/20 FU values differ."""
    if not matrix or not skus:
        return matrix
    powers = [_power_consumption_for_sku(s) for s in skus]
    if any(p is None for p in powers):
        return matrix
    out: list[list[str]] = []
    for row in matrix:
        if not row or row[0].strip() != "Потребляемая мощность":
            out.append(row)
            continue
        out.append(["Потребляемая мощность", *powers])  # type: ignore[list-item]
    return out


def _per_sku_values_by_nm(
    skus: list[str],
    by_nm: dict[int, str],
) -> list[str] | None:
    vals: list[str] = []
    for s in skus:
        nm = _sku_torque_nm(s)
        if nm is None or nm not in by_nm:
            return None
        vals.append(by_nm[nm])
    return vals


def _set_matrix_row(
    matrix: list[list[str]],
    label: str,
    values: list[str],
) -> list[list[str]]:
    out: list[list[str]] = []
    found = False
    for row in matrix:
        if row and row[0].strip() == label:
            out.append([label, *values])
            found = True
        else:
            out.append(row)
    if not found and values:
        out.append([label, *values])
    return out


def _normalize_fu_d_ds_function_rows(
    matrix: list[list[str]],
    skus: list[str],
) -> list[list[str]]:
    """Curate DA10/15/20 FU …-D/DS function rows from published PDFs."""
    if not matrix or not skus:
        return matrix
    if not all("FU" in s.upper() and re.search(r"-D(?:/DS)?$|-DS$", s, re.I) for s in skus):
        return matrix
    # Do not apply 10/15/20 shaft/mass/time overlays to DA3/DA5 FU.
    nms = {_sku_torque_nm(s) for s in skus}
    if not nms or not nms <= {10, 15, 20}:
        return matrix
    running = _per_sku_values_by_nm(skus, _RUNNING_FU_D_DS_BY_NM)
    mass = _per_sku_values_by_nm(skus, _MASS_FU_D_DS_BY_NM)
    if running:
        matrix = _set_matrix_row(matrix, "Время поворота", running)
    if mass:
        matrix = _set_matrix_row(matrix, "Масса", mass)
    n = len(skus)
    shared = {
        "Ручное управление": "отсутствует / не предусмотрено",
        "Температура хранения": "–40…+70 °C",
        "Испытание на влажность": "5…95 % относительной влажности",
        "Диаметр вала заслонки": "круглый 10…21 мм, квадратный 9×9…15×15 мм",
    }
    for label, val in shared.items():
        matrix = _set_matrix_row(matrix, label, [val] * n)
    return matrix


def _normalize_fu24_a_as_function_rows(
    matrix: list[list[str]],
    skus: list[str],
) -> list[list[str]]:
    """Curate DA10/15/20 FU24 …-A/AS function rows from published PDFs."""
    if not matrix or not skus:
        return matrix
    if not all(
        "FU" in s.upper() and "24" in s.upper() and re.search(r"-A(?:/AS)?$|-AS$", s, re.I)
        for s in skus
    ):
        return matrix
    nms = {_sku_torque_nm(s) for s in skus}
    if not nms or not nms <= {10, 15, 20}:
        return matrix
    running = _per_sku_values_by_nm(skus, _RUNNING_FU24_A_AS_BY_NM)
    area = _per_sku_values_by_nm(skus, _AREA_FU24_A_AS_BY_NM)
    mass = _per_sku_values_by_nm(skus, _MASS_FU24_A_AS_BY_NM)
    if running:
        matrix = _set_matrix_row(matrix, "Время поворота", running)
    if area:
        matrix = _set_matrix_row(matrix, "Площадь обслуживаемой заслонки", area)
        matrix = _set_matrix_row(matrix, "Макс. площадь заслонки", area)
    if mass:
        matrix = _set_matrix_row(matrix, "Масса", mass)
    n = len(skus)
    shared = {
        "Ручное управление": "отсутствует / не предусмотрено",
        "Температура хранения": "–40…+70 °C",
        "Испытание на влажность": "5…95 % относительной влажности",
        "Диаметр вала заслонки": "круглый 10…16 мм, квадратный 7×7…11×11 мм",
    }
    for label, val in shared.items():
        matrix = _set_matrix_row(matrix, label, [val] * n)
    # Dimensions ref: PPTX often says «чертежи»; PDF — «размеры».
    out: list[list[str]] = []
    for row in matrix:
        if row and row[0].strip().startswith("Габаритные размеры"):
            out.append([row[0], *["См. «Габаритные размеры»"] * max(len(row) - 1, n)])
        else:
            out.append(row)
    return out


def _normalize_da5fu_d_ds(doc: ManualDoc) -> None:
    """Align DA5FU D/DS lead + tech + aux with ``da5fu-d_ds.pdf``."""
    if doc.stem != "da5fu-d-ds":
        return
    if doc.summary:
        doc.summary = re.sub(
            r"Running\s*time\s*[：:]\s*<\s*70\s*с\s*,\s*"
            r"(?:время возврата пружины|Spring\s*reset(?:\s*time)?)\s*[：:]?\s*<\s*20\s*с",
            "Время поворота: < 70 с; время возврата пружины: < 20 с",
            doc.summary,
            flags=re.I,
        )
        doc.summary = re.sub(
            r"Исполнение\s*«S»\s*включает\s*2\s*вспомогательных?\s+переключателя\.?",
            "Исполнение «DS» включает 2 группы вспомогательных переключателей.",
            doc.summary,
            flags=re.I,
        )
        doc.summary = doc.summary.replace("Исполнение «S»", "Исполнение «DS»")

    # Aux factory table: cable terminals S1–S6 (not screw terminals 21–26).
    doc.aux_table = [
        ["Переключатель a", "Клеммы S1, S2", "Клеммы S1, S3"],
        ["–5…5°", "Замкнуто", "Разомкнуто"],
        ["5…90°", "Разомкнуто", "Замкнуто"],
        ["Переключатель b", "Клеммы S4, S5", "Клеммы S4, S6"],
        ["–5…80°", "Разомкнуто", "Замкнуто"],
        ["80…90°", "Замкнуто", "Разомкнуто"],
    ]

    n = max((_function_data_cols(doc.function_table) or 1), 1)
    doc.function_table = _set_matrix_row(
        doc.function_table,
        "Площадь обслуживаемой заслонки",
        ["< 0,5 м²"] * n,
    )
    # Alias label from some PPTX rows.
    doc.function_table = _set_matrix_row(
        doc.function_table,
        "Макс. площадь заслонки",
        ["< 0,5 м²"] * n,
    )
    doc.function_table = _set_matrix_row(
        doc.function_table,
        "Диаметр вала заслонки",
        ["круглый 10…16 мм, квадратный 7×7…11×11 мм"] * n,
    )
    doc.function_table = _set_matrix_row(
        doc.function_table,
        "Ручное управление",
        ["отсутствует / не предусмотрено"] * n,
    )
    doc.function_table = _set_matrix_row(
        doc.function_table,
        "Испытание на влажность",
        ["5…95 % относительной влажности"] * n,
    )
    # Wire sizing typo 0,5 м² → 0,5 мм² (PDF).
    if doc.electrical_table:
        fixed: list[list[str]] = []
        for row in doc.electrical_table:
            if row and "сечение" in row[0].lower():
                fixed.append(
                    [row[0], *[("0,5 мм²" if "м²" in c and "мм" not in c else c) for c in row[1:]]]
                )
            else:
                fixed.append(row)
        doc.electrical_table = fixed


def _normalize_da10_15_20fu24_a_as(doc: ManualDoc) -> None:
    """Align shared DA10/15/20 FU24 A/AS lead + aux + wire with published PDFs.

    Wiring/dims/rotation are identical across da10/15/20fu24-a:as.pdf.
    Factory-angle table: DA10/15 agree (0–10° / 10–90° / 0–80° / 80–90°);
    DA20 PDF differs (0–5° / 10–85° / 80–85°, S-labels). We follow DA10/15
    angles and S1–S6 labels from the aux schematic / faceplate.
    """
    if doc.stem != "da10-15-20fu24-a-as":
        return
    if doc.summary:
        doc.summary = re.sub(
            r"Исполнение\s*«S»\s*включает\s*2\s*вспомогательных?\s+переключателя\.?",
            "Исполнение «AS» включает 2 группы вспомогательных переключателей.",
            doc.summary,
            flags=re.I,
        )
        doc.summary = doc.summary.replace("Исполнение «S»", "Исполнение «AS»")

    # Match aux schematic / product cable labels (S1–S6); angles from DA10/15 PDF.
    doc.aux_table = [
        ["Переключатель a", "Клеммы S1, S2", "Клеммы S1, S3"],
        ["0–10°", "Замкнуто", "Разомкнуто"],
        ["10–90°", "Разомкнуто", "Замкнуто"],
        ["Переключатель b", "Клеммы S4, S5", "Клеммы S4, S6"],
        ["0–80°", "Разомкнуто", "Замкнуто"],
        ["80–90°", "Замкнуто", "Разомкнуто"],
    ]

    if doc.electrical_table:
        fixed: list[list[str]] = []
        for row in doc.electrical_table:
            if row and "сечение" in row[0].lower():
                fixed.append(
                    [row[0], *[("0,5 мм²" if "м²" in c and "мм" not in c else c) for c in row[1:]]]
                )
            else:
                fixed.append(row)
        doc.electrical_table = fixed


def _tech_sku_label(sku: str) -> str:
    """PDF-style label: ``DA4MU24-A`` → ``DA4MU24-A/AS``, ``…-D`` → ``…-D/DS``.

    SAFU keeps ``…-DS/DST``; SAMU stays ``…-DS`` (no DST in EN PDFs).
    """
    if "/DST" in sku.upper() or re.search(r"-DST$", sku, re.I):
        return sku
    if re.search(r"^SA\d*FU.*-DS$", sku, re.I):
        return re.sub(r"-DS$", "-DS/DST", sku, flags=re.I)
    if re.search(r"^SA\d*MU.*-DS$", sku, re.I):
        return sku
    if re.search(r"-AS$", sku, re.I):
        return re.sub(r"-AS$", "-A/AS", sku, flags=re.I)
    if re.search(r"-A$", sku, re.I):
        return f"{sku}/AS"
    if re.search(r"-DS$", sku, re.I):
        return re.sub(r"-DS$", "-D/DS", sku, flags=re.I)
    if re.search(r"-D$", sku, re.I):
        return f"{sku}/DS"
    return sku


def _tech_sku_header_label(sku: str) -> str:
    """ТТХ column header: keep edition (``…-A/AS`` / ``…-D/DS``), as in PDFs."""
    return _tech_sku_label(sku)


def _sku_header_row(skus: list[str], n_cols: int) -> list[str] | None:
    """Column headers under «Технические характеристики»."""
    if not skus or n_cols < 2:
        return None
    data_cols = n_cols - 1

    def join_labels(group: list[str]) -> str:
        labels: list[str] = []
        for s in group:
            label = _tech_sku_header_label(s)
            if label not in labels:
                labels.append(label)
        return "/".join(labels) if labels else "—"

    if data_cols == len(skus):
        labels: list[str] = []
        for s in skus:
            label = _tech_sku_header_label(s)
            if label not in labels:
                labels.append(label)
        # Per-SKU columns need one header per column (same width as data).
        if len(labels) == data_cols:
            return [""] + labels
        return [""] + [_tech_sku_header_label(s) for s in skus]
    v24, v230 = _sku_voltage_groups(skus)
    if data_cols == 2:
        return ["", join_labels(v24), join_labels(v230)]
    if data_cols == 1:
        return ["", join_labels(skus)]
    return None


def _function_data_cols(matrix: list[list[str]]) -> int | None:
    """Width of the richest data row in the function table (excl. label)."""
    if not matrix:
        return None
    best = 0
    for row in matrix:
        if not row:
            continue
        label = row[0].strip()
        rest = row[1:]
        if label and not any(c.strip() for c in rest):
            continue  # section banner
        filled = sum(1 for c in rest if c.strip())
        # Count filled only — PPTX often pads with blank ghost columns.
        best = max(best, filled)
    return best or None


def _compact_matrix_to_width(
    matrix: list[list[str]],
    n_data: int,
) -> list[list[str]]:
    """Keep label + non-empty values; pad/truncate to ``n_data`` columns."""
    if not matrix or n_data < 1:
        return matrix
    out: list[list[str]] = []
    for row in matrix:
        label = row[0] if row else ""
        rest = row[1:] if len(row) > 1 else []
        if label.strip() and not any(c.strip() for c in rest):
            out.append([label] + [""] * n_data)
            continue
        filled = [c for c in rest if c.strip()]
        if len(filled) == 1:
            out.append([label] + [filled[0]] * n_data)
        elif len(filled) >= n_data:
            out.append([label] + filled[:n_data])
        else:
            out.append([label] + filled + [""] * (n_data - len(filled)))
    return out


def _expand_voltage_cols_to_skus(
    matrix: list[list[str]],
    skus: list[str],
) -> list[list[str]]:
    """Expand label|24V|230V (or label|shared) rows into one column per SKU."""
    if not matrix or not skus:
        return matrix
    if len(matrix[0]) - 1 == len(skus):
        return matrix

    def src_idx(sku: str) -> int:
        return 1 if "230" in sku.upper() else 0

    out: list[list[str]] = []
    for row in matrix:
        label = row[0]
        rest = row[1:]
        if label.strip() and not any(c.strip() for c in rest):
            out.append([label] + [""] * len(skus))
            continue
        filled = [c for c in rest if c.strip()]
        if len(filled) == 1:
            out.append([label] + [filled[0]] * len(skus))
            continue
        if len(rest) == 2:
            out.append(
                [label]
                + [rest[src_idx(s)] if src_idx(s) < len(rest) else "" for s in skus]
            )
            continue
        out.append([label] + (rest + [""] * len(skus))[: len(skus)])
    return out


def _drop_empty_data_columns(matrix: list[list[str]]) -> list[list[str]]:
    """Remove data columns that are blank in every non-section row (PPTX ghosts)."""
    if not matrix:
        return matrix
    width = max(len(row) for row in matrix)
    if width < 2:
        return matrix
    keep = [True]
    for j in range(1, width):
        has_value = False
        for row in matrix:
            if len(row) <= j:
                continue
            label = row[0].strip()
            rest = row[1:]
            if label and not any(c.strip() for c in rest):
                continue
            if row[j].strip():
                has_value = True
                break
        keep.append(has_value)
    if all(keep):
        return matrix
    out: list[list[str]] = []
    for row in matrix:
        padded = row + [""] * (width - len(row))
        out.append([cell for cell, ok in zip(padded, keep) if ok])
    return out


def _slash_pair(left: str, right: str) -> str:
    """Join two cell values with « / »; factor a shared trailing unit when possible."""
    a, b = left.strip(), right.strip()
    if not a:
        return b
    if not b or a == b:
        return a
    for unit in (" м²", " кг", " °C", " мм", " Нм"):
        if a.endswith(unit) and b.endswith(unit):
            return f"{a[: -len(unit)].strip()} / {b[: -len(unit)].strip()}{unit}"
    return f"{a} / {b}"


def _merged_15_20_sku_code(sku15: str, sku20: str) -> str:
    """``DA15FU24-D`` + ``DA20FU24-D`` → ``DA15/20FU24-D``."""
    ma = re.match(r"^(DA)(\d+)(.*)$", sku15, re.I)
    mb = re.match(r"^(DA)(\d+)(.*)$", sku20, re.I)
    if ma and mb and ma.group(3).upper() == mb.group(3).upper():
        return f"{ma.group(1)}{ma.group(2)}/{mb.group(2)}{ma.group(3)}"
    return sku15


def _collapse_matrix_columns(
    matrix: list[list[str]],
    groups: list[list[int]],
    n_src: int,
) -> list[list[str]]:
    if not matrix:
        return matrix
    out: list[list[str]] = []
    for row in matrix:
        label = row[0] if row else ""
        rest = (row[1:] if len(row) > 1 else []) + [""] * n_src
        rest = rest[:n_src]
        if label.strip() and not any(c.strip() for c in rest):
            out.append([label] + [""] * len(groups))
            continue
        new_rest: list[str] = []
        for g in groups:
            if len(g) == 1:
                new_rest.append(rest[g[0]])
            else:
                new_rest.append(_slash_pair(rest[g[0]], rest[g[1]]))
        out.append([label] + new_rest)
    return out


def _collapse_nm_15_20_columns(
    elec: list[list[str]],
    func: list[list[str]],
    skus: list[str],
) -> tuple[list[list[str]], list[list[str]], list[str]]:
    """Merge adjacent DA15+DA20 columns (same voltage) into slash cells / DA15/20 headers."""
    if len(skus) < 2:
        return elec, func, skus
    groups: list[list[int]] = []
    i = 0
    while i < len(skus):
        nm = _sku_torque_nm(skus[i])
        volt = "230" if _sku_is_230v(skus[i]) else "24"
        if nm == 15 and i + 1 < len(skus):
            nm2 = _sku_torque_nm(skus[i + 1])
            volt2 = "230" if _sku_is_230v(skus[i + 1]) else "24"
            if nm2 == 20 and volt == volt2:
                groups.append([i, i + 1])
                i += 2
                continue
        groups.append([i])
        i += 1
    if all(len(g) == 1 for g in groups):
        return elec, func, skus

    new_skus: list[str] = []
    for g in groups:
        if len(g) == 1:
            new_skus.append(skus[g[0]])
        else:
            new_skus.append(_merged_15_20_sku_code(skus[g[0]], skus[g[1]]))

    n_src = len(skus)
    return (
        _collapse_matrix_columns(elec, groups, n_src),
        _collapse_matrix_columns(func, groups, n_src),
        new_skus,
    )


def _tech_tables_with_skus(
    electrical: list[list[str]],
    function: list[list[str]],
    skus: list[str],
) -> tuple[list[list[str]], list[list[str]]]:
    """Align electrical + function columns to SKUs; prepend SKU header on electrical."""
    if not skus:
        return electrical, function

    elec = _drop_empty_data_columns(electrical)
    func = _drop_empty_data_columns(function)

    per_sku = bool(skus) and _function_data_cols(func) == len(skus)
    if per_sku:
        elec = _expand_voltage_cols_to_skus(elec, skus)
        func = _compact_matrix_to_width(func, len(skus))
    elif func and _function_data_cols(func):
        # Still drop ghost blanks so colspan borders stay continuous.
        n = _function_data_cols(func) or 1
        func = _compact_matrix_to_width(func, n)

    func = _normalize_protection_class_rows(func, skus, per_sku=per_sku)
    func = _normalize_fu_d_ds_function_rows(func, skus)
    func = _normalize_fu24_a_as_function_rows(func, skus)
    elec = _normalize_power_consumption_rows(elec, skus)
    elec, func, skus = _collapse_nm_15_20_columns(elec, func, skus)

    if elec:
        header = _sku_header_row(skus, len(elec[0]))
        if header:
            elec = [header, *elec]
    return elec, func


def _torque_html(torque: str | None) -> str:
    """Render torque headline; 4+ values wrap mid-list into two lines."""
    text = (torque or "").strip()
    if not text:
        return ""
    parts = [p.strip() for p in text.split(" / ") if p.strip()]
    if len(parts) >= 4:
        mid = len(parts) // 2
        line1 = html.escape(" / ".join(parts[:mid]))
        line2 = html.escape(" / ".join(parts[mid:]))
        return f"{line1}<br>{line2}"
    return html.escape(text)


def _sku_block(skus: list[str]) -> str:
    if not skus:
        return ""
    v24, v230 = _sku_voltage_groups(skus)

    def unique_labels(group: list[str]) -> list[str]:
        labels: list[str] = []
        for s in group:
            label = _tech_sku_label(s)
            if label not in labels:
                labels.append(label)
        return labels

    # Two voltage columns: 24 V left, 230 V right (column-major fill).
    if v24 and v230:
        left = unique_labels(v24)
        right = unique_labels(v230)
        rows = max(len(left), len(right), 1)
        items = "".join(f"<li>{html.escape(label)}</li>" for label in [*left, *right])
        return (
            f'<ul class="sku-list sku-list-voltage" style="--sku-rows:{rows}">'
            f"{items}</ul>"
        )
    labels = unique_labels(list(skus))
    items = "".join(f"<li>{html.escape(label)}</li>" for label in labels)
    return f'<ul class="sku-list">{items}</ul>'


def _summary_parts(summary: str) -> tuple[str, str, str]:
    """Split summary into heading, intro prose, and list items from torque line on."""
    lines = [ln.strip() for ln in summary.split("\n") if ln.strip()]
    if not lines:
        return "", "", ""
    heading = html.escape(lines[0])
    body = lines[1:]
    intro = ""
    items: list[str] = []
    torque_re = re.compile(r"крутящий\s+момент\s*:", re.I)
    for i, ln in enumerate(body):
        m = torque_re.search(ln)
        if m:
            before = ln[: m.start()].strip()
            after = ln[m.start() :].strip()
            if before:
                intro = html.escape(before)
            if after:
                items.append(after)
            items.extend(body[i + 1 :])
            break
        if not items:
            # Keep narrative lines before specs as intro (join if several).
            intro = f"{intro} {html.escape(ln)}".strip() if intro else html.escape(ln)
    else:
        # No torque marker: treat remaining body as list if any.
        if body and not items:
            items = body
            intro = ""
    items_flat: list[str] = []
    for it in items:
        it = it.strip()
        if not it:
            continue
        # «Время поворота: …; время возврата пружины: …» → two bullets.
        m = re.match(
            r"(?is)^(Время поворота\s*:[^;]+);\s*(время возврата пружины\s*:.*)$",
            it,
        )
        if m:
            items_flat.append(m.group(1).strip())
            spring = m.group(2).strip()
            items_flat.append(spring[:1].upper() + spring[1:] if spring else spring)
            continue
        # «Номинальное напряжение: … Упр. сигнал Y: …» → two bullets.
        m = re.match(
            r"(?is)^(Номинальное напряжение\s*:.*)\s+(Упр\.\s*сигнал\s*Y\s*:.*)$",
            it,
        )
        if m:
            items_flat.append(m.group(1).strip())
            items_flat.append(m.group(2).strip())
            continue
        # «Обратная связь U: … Исполнение …» → two bullets.
        m = re.match(
            r"(?is)^(Обратная связь\s*U\s*:.*)\s+(Исполнение\s+.*)$",
            it,
        )
        if m:
            items_flat.append(m.group(1).strip())
            items_flat.append(m.group(2).strip())
            continue
        # SA: «Исполнения DS/DST …; DST — SAF72» → two bullets.
        m = re.match(
            r"(?is)^(Исполнения?\s+«DS»\s*/\s*«DST»[^;]*);?\s*"
            r"(«DST»\s*[—–-].*|Исполнение\s+«DST»\s*[—–-].*)$",
            it,
        )
        if m:
            items_flat.append(m.group(1).strip().rstrip(";").strip())
            dst = m.group(2).strip()
            if dst.startswith("«DST»"):
                dst = f"Исполнение {dst}"
            items_flat.append(dst)
            continue
        items_flat.append(it)
    items_html = "".join(f"<li>{html.escape(it)}</li>" for it in items_flat)
    list_html = f'<ul class="lead-list">{items_html}</ul>' if items_html else ""
    return heading, intro, list_html


def render_grid(doc: ManualDoc, *, diagram_profile: DiagramProfile | None = None) -> str:
    stem = doc.stem
    photo = (
        f'<img class="product-photo" src="assets/{stem}/{doc.product_photo}" alt="">'
        if doc.product_photo
        else '<div class="photo-fallback">HOOCON</div>'
    )
    if diagram_profile is not None:
        diagrams = curated_diagrams_html(stem, diagram_profile)
    else:
        diagrams = "".join(
            f'<figure class="diagram"><img src="assets/{stem}/{n}" alt=""></figure>'
            for n in doc.diagram_photos
        )
    lead = (
        f'<img class="lead-photo" src="assets/{stem}/{doc.lead_photo}" alt="">'
        if doc.lead_photo
        else ""
    )
    warnings = html.escape(doc.warnings).replace("\n", "<br>") if doc.warnings else "—"
    contacts_ru = html.escape(CONTACTS_HOOCON_RU).replace("\n", "<br>")
    contacts_by = html.escape(CONTACTS_CHAMPION_BY).replace("\n", "<br>")
    if doc.summary:
        summary_heading, summary_intro, summary_list = _summary_parts(doc.summary)
    else:
        summary_heading, summary_intro, summary_list = "", "", ""
    # Curated rotation figure already holds RU copy — skip EN PPTX footer note.
    rotation = ""
    if diagram_profile is None and doc.rotation_note:
        rotation = html.escape(doc.rotation_note)
    title = html.escape(doc.doc_title)
    title_parts = [p.strip() for p in re.split(r"\s+[—–-]\s+", doc.doc_title) if p.strip()]
    if len(title_parts) >= 2:
        title_html = "".join(
            f'<span class="doc-title-line">{html.escape(p)}</span>'
            for p in title_parts[:2]
        )
    else:
        title_html = f'<span class="doc-title-line">{title}</span>'

    tech_elec, tech_func = _tech_tables_with_skus(
        doc.electrical_table,
        doc.function_table,
        doc.skus,
    )

    aux_table_html = _html_table(doc.aux_table) or "<p class='muted'>—</p>"
    if diagram_profile is not None and diagram_profile.sheet1_aux_diagram:
        # Cols 1–6: aux table | diagram on top; Attention banner+prose full width below.
        dip_html = ""
        if diagram_profile.sheet1_dip_diagram:
            dip_html = f"""
            <div class="sheet1-dip-diagram">
              <img src="assets/{stem}/dip-diagram.png" alt="Режим управляющего сигнала">
            </div>"""
        aux_caption = diagram_profile.sheet1_aux_caption
        if aux_caption:
            sheet1_main = f"""
          <div class="stack sheet1-left sheet1-with-aux-diagram">
            <div class="sheet1-aux-row sheet1-aux-row-captioned">
              <div class="stack sheet1-aux sheet1-aux-main">
                <h2 class="banner">Вспомогательный переключатель</h2>
                {aux_table_html}
              </div>
              <div class="sheet1-aux-diagram">
                <div class="media aux-diagram-media aux-diagram-media-compact">
                  <img class="aux-diagram aux-diagram-compact" src="assets/{stem}/aux-diagram.png" alt="">
                </div>
              </div>
              <p class="aux-note">*Установите угол переключателя в соответствии с требованием заказчика</p>
              <p class="aux-diagram-caption">{html.escape(aux_caption)}</p>
            </div>{dip_html}
            <h2 class="banner banner-follow">Внимание</h2>
            <div class="prose">{warnings}</div>
          </div>
"""
        else:
            sheet1_main = f"""
          <div class="stack sheet1-left sheet1-with-aux-diagram">
            <div class="sheet1-aux-row">
              <div class="stack sheet1-aux">
                <h2 class="banner">Вспомогательный переключатель</h2>
                {aux_table_html}
                <p class="aux-note">*Установите угол переключателя в соответствии с требованием заказчика</p>
              </div>
              <div class="sheet1-aux-diagram">
                <div class="media aux-diagram-media">
                  <img class="aux-diagram" src="assets/{stem}/aux-diagram.png" alt="">
                </div>
              </div>
            </div>{dip_html}
            <h2 class="banner banner-follow">Внимание</h2>
            <div class="prose">{warnings}</div>
          </div>
"""
    else:
        sheet1_main = f"""
          <div class="stack sheet1-left">
            <h2 class="banner">Вспомогательный переключатель</h2>
            {aux_table_html}
            <p class="aux-note">*Установите угол переключателя в соответствии с требованием заказчика</p>
            <h2 class="banner banner-follow">Внимание</h2>
            <div class="prose">{warnings}</div>
          </div>
"""

    sheet1 = f"""
<section class="sheet sheet-page1" aria-label="Лист 1">
  <div class="col-guide" aria-hidden="true">{"".join(f"<span>{i}</span>" for i in range(1, 13))}</div>
  <div class="grid">
    <header class="span-12 hero">
      <p class="running-head">hoocon.ru</p>
      <div class="logo">
        <img src="assets/hoocon-logo.svg" alt="Hoocon" width="160" height="40">
      </div>
    </header>

    <div class="span-12 sheet1-body">
      <div class="sheet1-cols-left">
{sheet1_main}
        <div class="sheet1-contacts">
          <div class="stack contacts-ru">
            <h2 class="banner">ООО «ХОГОН»</h2>
            <div class="prose contacts">{contacts_ru}</div>
          </div>
          <div class="stack contacts-by">
            <h2 class="banner">ООО «Чемпион-Тэк»</h2>
            <div class="prose contacts">{contacts_by}</div>
          </div>
        </div>
      </div>
      <div class="sheet1-cols-right product-col">
        <div class="media">{photo}</div>
        <div class="media-meta">
          <h1 class="torque">{_torque_html(doc.torque)}</h1>
          {_sku_block(doc.skus)}
        </div>
        <div class="doc-title-box">
          <h2 class="doc-title">{title_html}</h2>
        </div>
      </div>
    </div>
    <p class="running-foot">hoocon.ru</p>
  </div>
</section>
"""

    sheet2 = f"""
<section class="sheet sheet-page2" aria-label="Лист 2">
  <div class="col-guide" aria-hidden="true">{"".join(f"<span>{i}</span>" for i in range(1, 13))}</div>
  <div class="grid">
    <header class="span-12 hero">
      <p class="running-head">hoocon.ru</p>
      <div class="logo">
        <img src="assets/hoocon-logo.svg" alt="Hoocon" width="160" height="40">
      </div>
    </header>

    <div class="span-12 sheet2-body">
      <div class="span-6 summary-block">
        <div class="summary-media">{lead or "<div class='photo-fallback'>HOOCON</div>"}</div>
        <div class="summary-copy">
          {f'<h2 class="lead-heading">{summary_heading}</h2>' if summary_heading else ""}
          {f'<p class="lead-intro">{summary_intro}</p>' if summary_intro else ""}
          {summary_list or '<p class="muted">—</p>'}
        </div>
      </div>

      <div class="span-6 stack tech-block">
        <h2 class="banner">Технические характеристики</h2>
        {_html_table(tech_elec, stretch_rows=True)}
        {_html_table(tech_func, stretch_rows=True)}
      </div>
      <div class="span-6 stack diagrams-block">
        <h2 class="banner">Схема подключения</h2>
        <div class="diagrams">{diagrams or "<p class='muted'>См. PDF</p>"}</div>
        {f'<p class="note">{rotation}</p>' if rotation else ""}
      </div>
    </div>
    <p class="running-foot">hoocon.ru</p>
  </div>
</section>
"""
    return sheet1 + sheet2


HTML_SHELL = """\
<!DOCTYPE html>
<html lang="ru">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{title}</title>
<style>
  :root {{
    --ink: #1a1a1a;
    --muted: #555;
    --line: #222;
    --banner: #231f20;
    --red: #d71519;
    --bg: #ececec;
    --paper: #fff;
    --gap: 3mm;
    --pad: 6mm;
  }}
  * {{ box-sizing: border-box; }}
  html, body {{
    margin: 0; padding: 0;
    background: var(--bg);
    color: var(--ink);
    font-family: "Helvetica Neue", Helvetica, Arial, "Noto Sans", sans-serif;
    font-size: 9pt;
    line-height: 1.35;
  }}
  .toolbar {{
    position: sticky; top: 0; z-index: 20;
    display: flex; flex-wrap: wrap; gap: 10px; align-items: center;
    padding: 10px 16px; background: #111; color: #fff; font-size: 13px;
  }}
  .toolbar a {{ color: #9cf; }}
  .toolbar button {{
    border: 0; border-radius: 4px; padding: 8px 14px;
    background: var(--red); color: #fff; font-weight: 600; cursor: pointer;
  }}
  .toolbar button.secondary {{
    background: #333;
  }}
  .toolbar button[aria-pressed="true"] {{
    outline: 2px solid #9cf;
    outline-offset: 1px;
  }}
  .stage {{
    padding: 12px; display: flex; flex-direction: column; gap: 12px; align-items: center;
  }}
  .sheet {{
    position: relative;
    display: flex;
    flex-direction: column;
    box-sizing: border-box;
    width: min(297mm, calc(100vw - 24px));
    aspect-ratio: 297 / 210;
    height: auto;
    max-height: calc(min(297mm, 100vw - 24px) * 210 / 297);
    overflow: hidden;
    background: var(--paper);
    box-shadow: 0 6px 20px rgba(0,0,0,.12);
    padding: var(--pad);
  }}
  .col-guide {{
    display: none;
    position: absolute;
    inset: var(--pad);
    z-index: 8;
    pointer-events: none;
    grid-template-columns: repeat(12, minmax(0, 1fr));
    gap: var(--gap);
  }}
  body.show-cols .col-guide {{ display: grid; }}
  .col-guide span {{
    min-height: 100%;
    background: rgba(215, 21, 25, 0.07);
    outline: 0.2mm dashed rgba(215, 21, 25, 0.4);
    color: rgba(215, 21, 25, 0.75);
    font-size: 8pt;
    font-weight: 700;
    text-align: center;
    padding-top: 1.2mm;
    letter-spacing: 0.02em;
  }}
  .col-guide span:nth-child(even) {{
    background: rgba(35, 100, 200, 0.07);
    outline-color: rgba(35, 100, 200, 0.4);
    color: rgba(35, 80, 160, 0.75);
  }}
  .grid {{
    position: relative;
    z-index: 1;
    display: grid;
    grid-template-columns: repeat(12, minmax(0, 1fr));
    gap: var(--gap);
    align-items: start;
    align-content: start;
    flex: 1 1 auto;
    width: 100%;
    min-height: 0;
    height: 100%;
  }}
  .template-voltage-note {{
    margin: 0 0 2mm;
    padding: 1.5mm 2.5mm;
    font-size: 7.5pt;
    line-height: 1.35;
    background: #f3f4f6;
    border-left: 2.5pt solid #b91c1c;
  }}
  .template-voltage-note code {{
    font-size: 0.92em;
  }}
  .sheet-page1 .grid {{
    grid-template-rows: auto minmax(0, 1fr) auto;
  }}
  .sheet-page1 .hero {{
    grid-row: 1;
    align-self: start;
  }}
  /* Sheet 1 body: two halves — cols 1–6 | 7–12. */
  .sheet-page1 .sheet1-body {{
    grid-row: 2;
    display: grid;
    grid-template-columns: repeat(12, minmax(0, 1fr));
    grid-template-rows: minmax(0, 1fr) auto;
    gap: var(--gap);
    align-self: stretch;
    min-height: 0;
    min-width: 0;
  }}
  .sheet-page1 .sheet1-cols-left {{
    grid-column: 1 / span 6;
    grid-row: 1 / span 2;
    display: flex;
    flex-direction: column;
    gap: var(--gap);
    min-height: 0;
    min-width: 0;
  }}
  .sheet-page1 .sheet1-cols-right {{
    grid-column: 7 / span 6;
    grid-row: 1;
    align-self: center;
    min-width: 0;
  }}
  .sheet-page1 .sheet1-left {{
    flex: 1 1 auto;
    display: flex;
    flex-direction: column;
    min-height: 0;
    min-width: 0;
  }}
  .sheet-page1 .sheet1-contacts {{
    flex: 0 0 auto;
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 1.5mm;
    align-items: start;
    min-width: 0;
  }}
  .sheet-page1 .sheet1-contacts .banner {{
    margin: 0 0 0.8mm;
    padding: 0.55mm 1.2mm;
    font-size: 7pt;
  }}
  .sheet-page1 .sheet1-contacts .prose.contacts {{
    font-size: 5.5pt;
    line-height: 1.15;
  }}
  .sheet-page1 .sheet1-aux {{
    /* Nested in .sheet1-aux-row — do not set grid-row (that stacked it under the diagram). */
    display: flex;
    flex-direction: column;
    min-height: 0;
    min-width: 0;
  }}
  .sheet-page1 .sheet1-left .prose,
  .sheet-page1 .sheet1-aux .prose {{
    flex: 1 1 auto;
  }}
  .sheet-page1 .sheet1-aux table {{
    font-size: 5.5pt;
    table-layout: fixed;
    width: 100%;
  }}
  .sheet-page1 .sheet1-aux th,
  .sheet-page1 .sheet1-aux td {{
    padding: 0.45mm 0.4mm;
    white-space: nowrap;
    overflow: hidden;
    text-overflow: clip;
  }}
  /* Aux table | diagram side-by-side; Attention below spans full 1–6. */
  .sheet-page1 .sheet1-aux-row {{
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: var(--gap);
    align-items: start;
    min-width: 0;
    flex: 0 1 auto;
    width: 100%;
  }}
  .sheet-page1 .sheet1-with-aux-diagram > .banner-follow {{
    margin-top: 1.5mm;
  }}
  .sheet-page1 .sheet1-with-aux-diagram > .prose {{
    /* Compact so DIP / aux diagrams keep room on sheet 1. */
    font-size: 6.2pt;
    line-height: 1.2;
    flex: 0 1 auto;
  }}
  .sheet-page1 .sheet1-aux-diagram {{
    min-width: 0;
    padding-top: 0;
    align-self: start;
  }}
  .sheet-page1 .aux-diagram-media {{
    min-height: 0;
    height: auto;
    max-height: 72mm;
    align-items: flex-start;
  }}
  .sheet-page1 .aux-diagram {{
    max-width: 100%;
    max-height: 72mm;
    width: 100%;
    height: auto;
    object-fit: contain;
    object-position: top center;
  }}
  .sheet-page1 .aux-diagram-media-compact {{
    max-height: 36mm;
  }}
  .sheet-page1 .aux-diagram-compact {{
    max-height: 36mm;
  }}
  .sheet-page1 .aux-diagram-caption {{
    margin: 0;
    padding: 0;
    font-size: 5.5pt;
    line-height: 1.25;
    color: #333;
  }}
  /* Captioned aux: note + caption share one grid row (same baseline). */
  .sheet-page1 .sheet1-aux-row-captioned {{
    grid-template-rows: minmax(0, 1fr) auto;
    align-items: end;
  }}
  .sheet-page1 .sheet1-aux-row-captioned .sheet1-aux-main {{
    grid-column: 1;
    grid-row: 1;
    align-self: stretch;
  }}
  .sheet-page1 .sheet1-aux-row-captioned .sheet1-aux-diagram {{
    grid-column: 2;
    grid-row: 1;
    align-self: end;
    justify-self: stretch;
  }}
  .sheet-page1 .sheet1-aux-row-captioned > .aux-note {{
    grid-column: 1;
    grid-row: 2;
    margin: 1mm 0 0;
  }}
  .sheet-page1 .sheet1-aux-row-captioned > .aux-diagram-caption {{
    grid-column: 2;
    grid-row: 2;
    margin: 1mm 0 0;
  }}
  .sheet-page1 .sheet1-dip-diagram {{
    min-height: 0;
    max-height: 48mm;
    overflow: hidden;
    flex: 0 0 auto;
  }}
  .sheet-page1 .sheet1-dip-diagram img {{
    display: block;
    width: 100%;
    max-height: 48mm;
    height: auto;
    object-fit: contain;
    object-position: left center;
  }}
  .sheet-page1 .running-foot {{
    grid-row: 3;
  }}
  .sheet-page2 .grid {{
    grid-template-rows: auto minmax(0, 1fr) auto;
  }}
  .sheet-page2 .hero {{
    grid-row: 1;
    align-self: start;
  }}
  .sheet-page2 .sheet2-body {{
    grid-row: 2;
    display: grid;
    grid-template-columns: repeat(12, minmax(0, 1fr));
    grid-template-rows: auto minmax(0, 1fr);
    gap: var(--gap);
    align-content: stretch;
    align-self: stretch;
    min-height: 0;
    min-width: 0;
    padding-bottom: 1mm;
  }}
  .sheet-page2 .running-foot {{
    grid-row: 3;
  }}
  .sheet2-body .summary-block {{
    grid-column: 1 / span 6;
    grid-row: 1;
    align-self: start;
  }}
  .sheet2-body .tech-block {{
    grid-column: 1 / span 6;
    grid-row: 2;
    align-self: stretch;
    min-height: 0;
    height: 100%;
    display: flex;
    flex-direction: column;
    overflow: hidden;
    /* Keep collapsed bottom border inside the clip box. */
    padding-bottom: 0.35mm;
  }}
  .sheet2-body .tech-block .banner {{
    flex: 0 0 auto;
  }}
  /* Div slots shrink under flex; bare <table> items often will not. */
  .sheet2-body .tech-block .tech-table-slot {{
    min-height: 0;
    overflow: hidden;
    display: flex;
    flex-direction: column;
    margin: 0 0 1.2mm;
    position: relative;
    /* Room for the collapsed outer bottom border (overflow:hidden clips it). */
    padding-bottom: 2px;
  }}
  /* Electrical (first) slot: size to content so wrapped «Потребляемая мощность»
     does not clip the following «Сечение…» row. */
  .sheet2-body .tech-block .tech-table-slot:not(:last-child) {{
    flex: 0 0 auto;
    overflow: visible;
    min-height: auto;
  }}
  .sheet2-body .tech-block .tech-table-slot:not(:last-child) .data-table,
  .sheet2-body .tech-block .tech-table-slot:not(:last-child) .data-table tbody,
  .sheet2-body .tech-block .tech-table-slot:not(:last-child) .data-table tr {{
    height: auto;
  }}
  .sheet2-body .tech-block .tech-table-slot:last-child {{
    margin-bottom: 0;
  }}
  .sheet2-body .tech-block .tech-table-slot:last-child::after {{
    /* Always-visible bottom edge when cell borders are subpixel-clipped. */
    content: "";
    position: absolute;
    left: 0;
    right: 0;
    bottom: 0;
    height: 0.25mm;
    background: #444;
    pointer-events: none;
  }}
  .sheet2-body .tech-block .data-table {{
    flex: 1 1 auto;
    width: 100%;
    height: 100%;
    margin: 0;
    table-layout: fixed;
    border-collapse: collapse;
    font-size: 5.2pt;
  }}
  .sheet2-body .tech-block .data-table tbody {{
    height: 100%;
  }}
  .sheet2-body .tech-block .data-table tr {{
    height: calc(100% / var(--tech-rows, 1));
  }}
  .sheet2-body .tech-block .data-table td {{
    text-align: center;
    vertical-align: middle;
    /* ~1–2px vertical gap; horizontal stays readable. */
    padding: 1px 0.8mm;
    line-height: 1.1;
    box-sizing: border-box;
    border: 0.2mm solid #444;
  }}
  .sheet2-body .tech-block .data-table tr:last-child td {{
    border-bottom: 0.25mm solid #444;
  }}
  .sheet2-body .tech-block .data-table td:first-child,
  .sheet2-body .tech-block .data-table td.row-section {{
    text-align: left;
  }}
  .sheet2-body .tech-block .data-table tr:first-child td {{
    background: #e8e8e8;
    font-weight: 650;
    text-align: center;
    font-size: 5pt;
    white-space: nowrap;
    word-break: normal;
    overflow: hidden;
    text-overflow: clip;
  }}
  .sheet2-body .tech-block .data-table tr:first-child td:first-child {{
    text-align: left;
  }}
  .sheet2-body .tech-block .data-table tr:first-child td:not(:first-child) {{
    text-align: center;
  }}
  .sheet2-body .diagrams-block {{
    grid-column: 7 / span 6;
    grid-row: 1 / span 2;
    align-self: stretch;
    min-height: 0;
    height: 100%;
    display: flex;
    flex-direction: column;
    overflow: hidden;
  }}
  .sheet2-body .diagrams-block .diagrams {{
    flex: 1 1 auto;
    min-height: 0;
    display: flex;
    flex-direction: column;
    gap: 1.5mm;
    overflow: hidden;
  }}
  .sheet2-body .diagrams-block .diagrams > .diagram {{
    flex: 0 0 auto;
    min-height: 0;
  }}
  .sheet2-body .diagrams-block .diagrams > .diagram-dip {{
    flex: 1 1 auto;
    min-height: 0;
    overflow: hidden;
    display: flex;
    flex-direction: column;
    align-items: stretch;
  }}
  .sheet-page2 .lead-heading {{
    font-size: clamp(10pt, 1.8vw, 13pt);
    margin-bottom: 1mm;
  }}
  .sheet-page2 .lead-intro {{
    font-size: 7pt;
    margin-bottom: 1mm;
    line-height: 1.25;
  }}
  .sheet-page2 .lead-list {{
    font-size: 7pt;
    padding-left: 3.2mm;
    line-height: 1.25;
  }}
  .sheet-page2 .lead-list li {{
    margin: 0 0 0.4mm;
  }}
  .sheet-page2 .lead-photo {{
    max-height: 28mm;
  }}
  .sheet-page2 .banner {{
    margin: 0 0 1.2mm;
    padding: 0.8mm 1.5mm;
    font-size: 7.5pt;
  }}
  .sheet-page2 .data-table {{
    margin: 0 0 1.5mm;
    font-size: 5.8pt;
    line-height: 1.2;
  }}
  .sheet-page2 .data-table td {{
    padding: 0.45mm 0.7mm;
  }}
  .sheet-page2 .diagram {{
    min-height: 0;
    padding: 0.8mm;
  }}
  .sheet-page2 .diagrams {{
    gap: 1.5mm;
    grid-template-columns: 1fr;
  }}
  .sheet-page2 .diagram-wide img {{
    width: 100%;
    max-height: 34mm;
  }}
  .sheet-page2 .wiring-board > img {{
    max-height: 40mm;
  }}
  .wiring-board {{
    position: relative;
    width: 100%;
  }}
  .wiring-ru-headers {{
    display: grid;
    grid-template-columns: 1.55fr 1fr;
    gap: 1mm 2mm;
    margin: 0 0 1mm;
    font-size: 7pt;
    font-weight: 650;
    line-height: 1.15;
    text-align: center;
    color: #111;
  }}
  .wiring-ru-headers .wrh-act {{
    /* spans both actuator voltage blocks */
  }}
  .wiring-board > img {{
    display: block;
    width: 100%;
    height: auto;
    max-height: 40mm;
    object-fit: contain;
  }}
  .wiring-board-ru-headers > img {{
    max-height: 40mm;
  }}
  .wiring-board .wl {{
    position: absolute;
    font-size: 5.5pt;
    font-weight: 650;
    line-height: 1.1;
    color: #111;
    white-space: nowrap;
    pointer-events: none;
  }}
  .wl-p2a {{ left: 3%; top: 1%; }}
  .wl-p3a {{ left: 18%; top: 1%; }}
  .wl-p2b {{ left: 35.5%; top: 0%; }}
  .wl-p3b {{ left: 50.5%; top: 0%; }}
  .wl-v24 {{ left: 6%; top: 88%; }}
  .wl-v230 {{ left: 38%; top: 88%; }}
  .wl-rating {{ left: 76%; top: 88%; }}
  .sheet-page2 .diagrams .diagram:nth-child(2) img {{
    max-height: 48mm;
  }}
  .sheet-page2 .diagrams .diagram:nth-child(3) img {{
    max-height: 42mm;
  }}
  .sheet-page2 .diagrams .diagram-dip > img {{
    max-height: none;
  }}
  /* Free vertical room for the wide DIP table on A/AS sheet 2. */
  .sheet-page2 .diagrams:has(.diagram-dip) .wiring-board > img {{
    max-height: 28mm !important;
  }}
  .sheet-page2 .diagrams:has(.diagram-dip) .diagram:nth-child(2) img {{
    max-height: 38mm !important;
  }}
  .diagrams .banner {{
    display: block;
    width: 100%;
    box-sizing: border-box;
  }}
  .diagram-rotation .rotation-panel {{
    display: grid;
    /* diagrams-block = sheet cols 7–12 → local 6; image on sheet 10–12 = local 4–6. */
    grid-template-columns: repeat(6, minmax(0, 1fr));
    gap: var(--gap);
    align-items: start;
    width: 100%;
    min-height: 0;
  }}
  .diagram-rotation .rotation-panel .rotation-copy {{
    grid-column: 1 / span 3;
    width: 100%;
    box-sizing: border-box;
    padding: 0;
    min-width: 0;
  }}
  .diagram-rotation .rotation-panel > img,
  .diagram-rotation .rotation-panel > .rotation-media {{
    grid-column: 4 / span 3;
    width: 100%;
    min-width: 0;
    justify-self: stretch;
  }}
  .diagram-rotation .rotation-media {{
    display: flex;
    flex-direction: column;
    align-items: stretch;
    gap: 0.6mm;
    min-height: 0;
  }}
  .diagram-rotation .rotation-media > img {{
    width: 100%;
    max-height: 34mm !important;
    object-fit: contain;
    object-position: top center;
  }}
  .diagram-rotation .rotation-media > .rotation-label {{
    margin: 0 !important;
    text-align: center;
    font-size: 6pt;
  }}
  .diagram-rotation > .rotation-copy {{
    width: 100%;
    box-sizing: border-box;
    padding: 0.5mm 0;
  }}
  .rotation-copy {{
    font-size: 6pt;
    line-height: 1.25;
  }}
  .rotation-copy p {{
    margin: 0 0 0.8mm;
  }}
  .rotation-table {{
    margin: 0 0 0.8mm !important;
    font-size: 5.2pt !important;
    table-layout: fixed;
    width: 100%;
  }}
  .rotation-table td {{
    white-space: nowrap;
    padding: 0.3mm 0.5mm !important;
    overflow: hidden;
    text-overflow: clip;
  }}
  .rotation-table td:first-child {{
    width: 36%;
  }}
  .rotation-label {{
    margin: 0 !important;
    font-weight: 650;
  }}
  .diagram-rotation img {{
    max-height: 36mm !important;
  }}
  .diagram-rotation.diagram-dip > img {{
    width: 100%;
    flex: 1 1 auto;
    min-height: 0;
    /* Definite height so object-fit can scale the whole table (not clip). */
    height: 0;
    max-height: none !important;
    object-fit: contain;
    object-position: top center;
    display: block;
  }}
  .diagram-rotation.diagram-thermal .rotation-panel > img {{
    max-height: 52mm !important;
    width: 100%;
    height: auto;
    object-fit: contain;
    object-position: top center;
  }}
  .sheet-page2 .note {{
    margin: 1mm 0 0;
    padding: 1mm 1.5mm;
    font-size: 6.5pt;
  }}
  .span-3 {{ grid-column: span 3; }}
  .span-4 {{ grid-column: span 4; }}
  .span-5 {{ grid-column: span 5; }}
  .span-6 {{ grid-column: span 6; }}
  .span-7 {{ grid-column: span 7; }}
  .span-12 {{ grid-column: span 12; }}
  .stack {{ min-width: 0; }}
  .summary-block {{
    grid-column: 1 / span 6;
    display: grid;
    grid-template-columns: repeat(6, minmax(0, 1fr));
    gap: var(--gap);
    align-items: center;
    min-width: 0;
  }}
  .summary-media {{
    grid-column: 1 / span 2;
    display: flex;
    align-items: center;
    justify-content: center;
    min-width: 0;
  }}
  .lead-photo {{
    display: block;
    max-width: 100%;
    max-height: 42mm;
    object-fit: contain;
  }}
  .summary-copy {{
    grid-column: 3 / span 4;
    min-width: 0;
  }}
  .lead-heading {{
    margin: 0 0 1.5mm;
    font-size: clamp(12pt, 2.2vw, 16pt);
    font-weight: 700;
    line-height: 1.2;
  }}
  .lead-intro {{
    margin: 0 0 2mm;
    font-size: 8.5pt;
  }}
  .lead-list {{
    margin: 0;
    padding-left: 4mm;
    font-size: 8.5pt;
  }}
  .lead-list li {{
    margin: 0 0 0.8mm;
  }}

  .hero {{
    display: grid;
    grid-template-columns: repeat(12, minmax(0, 1fr));
    gap: var(--gap);
    align-items: center;
    padding-bottom: 2.5mm;
  }}
  .running-head {{
    grid-column: 1 / span 1;
    margin: 0;
    font-size: 8pt;
    font-weight: 600;
    letter-spacing: 0.02em;
    color: var(--muted);
    white-space: nowrap;
  }}
  .running-foot {{
    grid-column: 12 / span 1;
    margin: 0;
    font-size: 8pt;
    font-weight: 600;
    letter-spacing: 0.02em;
    color: var(--muted);
    white-space: nowrap;
    text-align: right;
    justify-self: end;
    align-self: end;
  }}
  .logo {{
    grid-column: 11 / span 2;
    display: flex;
    align-items: center;
    justify-self: end;
    flex-shrink: 0;
  }}
  .logo img {{
    display: block;
    height: 9mm;
    width: auto;
    max-width: 100%;
  }}
  .product-col {{
    display: grid;
    grid-template-columns: repeat(6, minmax(0, 1fr));
    gap: var(--gap);
    align-items: center;
    min-width: 0;
  }}
  .product-col .media {{
    /* Sheet cols 7–9 (product-col is sheet cols 7–12 → local 1 / span 3). */
    grid-column: 1 / span 3;
    grid-row: 1;
    align-self: center;
    justify-self: center;
    min-height: 0;
    width: 100%;
  }}
  .sheet-page1 .product-col .product-photo {{
    max-width: 100%;
    max-height: 72mm;
    width: auto;
    height: auto;
    object-fit: contain;
  }}
  /* SA FU/MU: slightly smaller hero so left body + SAF72 fit without clip. */
  .sheet-page1 .product-col .product-photo[src*="/sa"] {{
    max-height: 58mm;
  }}
  .doc-title-box {{
    grid-column: 2 / span 5;
    grid-row: 2;
    container-type: inline-size;
    container-name: doc-title;
    width: 100%;
    min-width: 0;
    align-self: center;
  }}
  .doc-title {{
    display: flex;
    flex-direction: column;
    justify-content: center;
    margin: 0;
    width: 100%;
    overflow: visible;
    font-size: clamp(9pt, 5.8cqw, 16pt);
    font-weight: 700;
    line-height: 1.2;
    text-align: right;
  }}
  .doc-title-line {{
    display: block;
    width: 100%;
    box-sizing: border-box;
    white-space: nowrap;
    overflow: visible;
    text-overflow: clip;
  }}
  .media-meta {{
    grid-column: 4 / span 3;
    grid-row: 1;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    gap: 1.5mm;
    align-self: center;
    justify-self: center;
    min-width: 0;
    width: 100%;
    text-align: center;
    container-type: inline-size;
    container-name: media-meta;
  }}
  .torque {{
    margin: 0;
    width: 100%;
    max-width: 100%;
    box-sizing: border-box;
    font-size: clamp(14pt, 16cqw, 48pt);
    font-weight: 700;
    color: var(--red);
    line-height: 1.1;
    white-space: nowrap;
    text-align: center;
    overflow: hidden;
  }}
  .sku-list {{
    list-style: none;
    margin: 0;
    padding: 2mm;
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 0.2mm 1mm;
    justify-items: center;
    justify-content: center;
    align-self: stretch;
    width: 100%;
    max-width: 100%;
    font-size: 8pt;
    font-weight: 700;
    font-variant-numeric: tabular-nums;
    background: var(--banner);
    color: #fff;
    box-sizing: border-box;
    text-align: center;
  }}
  /* Mixed 24 V + 230 V: left column = 24 V, right = 230 V (column-major). */
  .sku-list-voltage {{
    grid-template-columns: 1fr 1fr;
    grid-auto-flow: column;
    grid-template-rows: repeat(var(--sku-rows, 1), auto);
  }}
  .sku-list li {{
    white-space: nowrap;
    text-align: center;
  }}

  .banner {{
    margin: 0 0 2mm; padding: 1.2mm 2mm;
    background: var(--banner); color: #fff;
    font-size: 8.5pt; font-weight: 650;
  }}
  .banner-follow {{
    margin-top: 3.5mm;
  }}
  .aux-note {{
    margin: 1.5mm 0 0;
    font-size: 7.5pt;
    color: var(--muted);
    line-height: 1.3;
  }}
  .attention-block .prose {{
    margin-bottom: 0;
  }}
  .prose {{ font-size: 8pt; }}
  .prose.contacts {{ font-size: 7.5pt; color: var(--muted); }}
  .lead-text {{ font-size: 8.5pt; }}
  .muted {{ color: var(--muted); }}
  .note {{
    margin: 2mm 0 0; padding: 1.5mm 2mm;
    border-left: 0.55mm solid var(--red); font-size: 8pt;
  }}
  .media {{
    display: flex; align-items: center; justify-content: center;
    min-height: 48mm; border: none; background: transparent;
  }}
  .product-photo, .photo-fallback {{
    max-width: 100%; max-height: 62mm; object-fit: contain;
  }}
  .photo-fallback {{
    display: grid; place-items: center; width: 100%; height: 48mm;
    color: #bbb; font-weight: 700; letter-spacing: .12em;
    background: transparent;
  }}
  .diagrams {{
    display: grid; grid-template-columns: 1fr 1fr; gap: 2mm;
  }}
  .diagram {{
    margin: 0; border: none; background: transparent;
    display: flex; align-items: center; justify-content: center;
    min-height: 30mm; padding: 1.5mm;
  }}
  .diagram:has(> .banner) {{
    flex-direction: column;
    align-items: stretch;
    justify-content: flex-start;
  }}
  .diagram > .banner {{
    display: block;
    width: 100%;
    box-sizing: border-box;
    flex: 0 0 auto;
  }}
  .diagram-wide {{
    grid-column: 1 / -1;
  }}
  .diagram img {{ max-width: 100%; max-height: 44mm; object-fit: contain; }}
  .diagram-wide img {{
    max-height: 36mm;
    width: 100%;
    object-fit: contain;
  }}
  .data-table {{
    width: 100%; border-collapse: collapse; table-layout: fixed;
    margin: 0 0 2.5mm; font-size: 7.2pt;
  }}
  .data-table td {{
    border: 0.2mm solid #444; padding: 1mm 1.1mm;
    vertical-align: middle; text-align: center; word-break: break-word;
  }}
  .data-table tr:first-child td,
  .data-table td.row-section {{
    background: #e8e8e8; font-weight: 600; text-align: left;
  }}

  @media (max-width: 900px) {{
    .span-3, .span-4, .span-5, .span-6, .span-7 {{ grid-column: span 12; }}
    .summary-block {{ grid-column: span 12; grid-template-columns: 1fr; }}
    .summary-media,
    .summary-copy {{ grid-column: 1 / -1; }}
    .sheet2-body .summary-block,
    .sheet2-body .tech-block,
    .sheet2-body .diagrams-block {{
      grid-column: 1 / -1;
      grid-row: auto;
    }}
    .sheet1-body .sheet1-cols-left,
    .sheet1-body .sheet1-cols-right {{
      grid-column: 1 / -1;
      grid-row: auto;
    }}
    .product-col {{ grid-template-columns: 1fr; }}
    .product-col .media,
    .media-meta,
    .doc-title-box,
    .torque,
    .sku-list {{ grid-column: 1 / -1; }}
    .diagrams {{ grid-template-columns: 1fr; }}
    .diagram-rotation .rotation-panel {{ grid-template-columns: 1fr; }}
    .diagram-rotation .rotation-panel .rotation-copy,
    .diagram-rotation .rotation-panel > img,
    .diagram-rotation .rotation-panel > .rotation-media {{ grid-column: 1 / -1; }}
  }}

  @page {{ size: A4 landscape; margin: 0; }}
  @media print {{
    html, body {{ background: #fff; }}
    .toolbar {{ display: none !important; }}
    .col-guide {{ display: none !important; }}
    .stage {{ padding: 0; gap: 0; }}
    .sheet {{
      box-shadow: none;
      width: 297mm;
      height: 210mm;
      max-height: 210mm;
      aspect-ratio: 297 / 210;
      page-break-after: always;
      break-after: page;
      overflow: hidden;
    }}
    .sheet:last-child {{ page-break-after: auto; }}
    .span-3 {{ grid-column: span 3; }}
    .span-4 {{ grid-column: span 4; }}
    .span-5 {{ grid-column: span 5; }}
    .span-6 {{ grid-column: span 6; }}
    .span-7 {{ grid-column: span 7; }}
    .summary-block {{
      grid-column: 1 / span 6;
      grid-template-columns: repeat(6, minmax(0, 1fr));
    }}
    .summary-media {{ grid-column: 1 / span 2; }}
    .summary-copy {{ grid-column: 3 / span 4; }}
    /* Photo crop (clip-path) must remain for print PDF — only hide the UI. */
    .photo-crop-panel,
    #toggle-photo-crop {{ display: none !important; }}
  }}
</style>
</head>
<body class="show-cols">
  <div class="toolbar">
    <strong>{title}</strong>
    <span>2 листа / A4 альбом / сетка 12</span>
    <button type="button" id="toggle-cols" class="secondary" aria-pressed="true">Колонки</button>
    <button type="button" onclick="window.print()">Печать / PDF</button>
    <a href="index.html">все инструкции</a>
  </div>
  <div class="stage">
{body}
  </div>
  <script>
    (function () {{
      var btn = document.getElementById("toggle-cols");
      if (btn) {{
        btn.addEventListener("click", function () {{
          var on = document.body.classList.toggle("show-cols");
          btn.setAttribute("aria-pressed", on ? "true" : "false");
          window.requestAnimationFrame(fitSheet1Type);
        }});
      }}

      /* Shrink title lines to box width — never ellipsis / clip. Shared size. */
      function fitDocTitle() {{
        document.querySelectorAll(".doc-title").forEach(function (title) {{
          var box = title.closest(".doc-title-box") || title;
          var boxW = box.clientWidth;
          if (boxW < 8) return;
          var lines = title.querySelectorAll(".doc-title-line");
          if (!lines.length) lines = [title];
          var lo = 8;
          var hi = 18;
          var best = 8;
          for (var i = 0; i < 18; i++) {{
            var mid = (lo + hi) / 2;
            for (var k = 0; k < lines.length; k++) {{
              lines[k].style.fontSize = mid + "pt";
            }}
            var overflow = false;
            for (var n = 0; n < lines.length; n++) {{
              if (lines[n].scrollWidth > boxW + 0.5) {{
                overflow = true;
                break;
              }}
            }}
            if (overflow) hi = mid;
            else {{
              best = mid;
              lo = mid;
            }}
          }}
          for (var m = 0; m < lines.length; m++) {{
            lines[m].style.fontSize = best + "pt";
          }}
        }});
      }}

      /* Shrink torque headline to media-meta width (e.g. «10 / 15 / 20 Нм»). */
      function fitTorque() {{
        document.querySelectorAll(".torque").forEach(function (el) {{
          var box = el.closest(".media-meta") || el.parentElement;
          var boxW = box ? box.clientWidth : 0;
          if (boxW < 8) return;
          var lo = 12;
          var hi = 48;
          var best = 12;
          for (var i = 0; i < 20; i++) {{
            var mid = (lo + hi) / 2;
            el.style.fontSize = mid + "pt";
            if (el.scrollWidth > boxW + 0.5) hi = mid;
            else {{
              best = mid;
              lo = mid;
            }}
          }}
          el.style.fontSize = best + "pt";
        }});
      }}

      function fitSheet1Type() {{
        fitDocTitle();
        fitTorque();
      }}

      window.addEventListener("load", fitSheet1Type);
      window.addEventListener("resize", fitSheet1Type);
      window.addEventListener("beforeprint", fitSheet1Type);
      if (document.fonts && document.fonts.ready) {{
        document.fonts.ready.then(fitSheet1Type);
      }}
    }})();
  </script>
  <script src="assets/photo-crop-tool.js"></script>
</body>
</html>
"""

STEM_MAP = {
    "DA2MU24(230)-D(S).pptx": ("da2mu-d-ds", "DA2MU …-D/DS — руководство (RU)"),
    "DA2MU24(230)-A(S).pptx": ("da2mu-a-as", "DA2MU …-A/AS — руководство (RU)"),
    "DA3FU24(230)-D(S).pptx": ("da3fu-d-ds", "DA3FU …-D/DS — руководство (RU)"),
    "DA4_6MU24(230)-A(S).pptx": ("da4-6mu-a-as", "DA4/6MU …-A/AS — руководство (RU)"),
    "DA4_6MU24(230)-D(S).pptx": ("da4-6mu-d-ds", "DA4/6MU …-D/DS — руководство (RU)"),
    "DA5FU24(230)-D(S).pptx": ("da5fu-d-ds", "DA5FU …-D/DS — руководство (RU)"),
    "DA10_15_20FU24-A(S).pptx": (
        "da10-15-20fu24-a-as",
        "DA10/15/20FU24 …-A/AS — руководство (RU)",
    ),
    "DA10_15_20FU24(230)-D(S).pptx": (
        "da10-15-20fu24-230-d-ds",
        "DA10/15/20FU24/230 …-D/DS — руководство (RU)",
    ),
}


def _set_runs_text(shape_or_cell, new_text: str) -> None:
    tf = shape_or_cell.text_frame
    paragraphs = list(tf.paragraphs)
    lines = new_text.split("\n")
    for i, line in enumerate(lines):
        if i < len(paragraphs):
            p = paragraphs[i]
            if p.runs:
                p.runs[0].text = line
                for run in p.runs[1:]:
                    run.text = ""
            else:
                p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    for j in range(len(lines), len(paragraphs)):
        p = paragraphs[j]
        if p.runs:
            p.runs[0].text = ""
            for run in p.runs[1:]:
                run.text = ""
        else:
            p.text = ""


def translate_pptx_inplace(prs: Presentation) -> None:
    for slide in prs.slides:
        for sh, _l, _t in iter_shapes(slide.shapes):
            if sh.has_text_frame:
                raw = "\n".join(p.text for p in sh.text_frame.paragraphs)
                if raw.strip():
                    _set_runs_text(sh, translate(raw))
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        raw = cell.text
                        if raw.strip():
                            _set_runs_text(cell, translate(raw))


def convert_one(src: Path, out_dir: Path, *, force: bool = False) -> Path:
    meta = STEM_MAP.get(src.name)
    if not meta:
        raise SystemExit(f"Unknown PPTX mapping: {src.name}")
    stem, title = meta
    finished_dir = finished_manuals_dir(out_dir, stem=stem)
    finished_dir.mkdir(parents=True, exist_ok=True)
    _ensure_family_logo(finished_dir)
    if manual_stem_is_locked(stem) and not force:
        html_path = finished_dir / f"{stem}.html"
        if html_path.is_file():
            return html_path
        raise SystemExit(
            f"Locked manual {stem!r} has no HTML yet; pass force=True to build."
        )
    assets = finished_dir / "assets" / stem
    if assets.exists():
        for old in assets.iterdir():
            if old.is_file():
                old.unlink()
    assets.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(src))
    translate_pptx_inplace(prs)
    # PDF instructions are canon; do not publish a translated PPTX beside HTML.
    doc = extract_manual(prs, assets, stem, title)
    diagram_profile = ensure_diagram_assets(stem, finished_dir, force=force)
    # Prefer photos cropped from the instruction PDF when present.
    if (assets / "product.png").is_file():
        doc.product_photo = "product.png"
        doc.lead_photo = (
            "lead.png" if (assets / "lead.png").is_file() else "product.png"
        )
    body = render_grid(doc, diagram_profile=diagram_profile)
    html_path = finished_dir / f"{stem}.html"
    html_path.write_text(
        HTML_SHELL.format(title=html.escape(title), body=body),
        encoding="utf-8",
    )
    return html_path


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--src-dir",
        type=Path,
        default=Path("/Users/niko/Downloads/Док-ция"),
    )
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=_MANUALS_RU_DIR,
    )
    args = parser.parse_args()
    args.out_dir.mkdir(parents=True, exist_ok=True)

    for path in write_voltage_template_shells(args.out_dir):
        print("OK", path.name)

    index_rows: list[str] = [
        "<li><strong>Шаблоны напряжения (для EN→RU)</strong><br>"
        '<a href="template-v24.html">V24 — AC/DC 24 В</a> · '
        '<a href="template-v230.html">V230 — AC 100…240 В</a> · '
        "<code>TEMPLATES.md</code></li>",
        "<li><strong>Семейства</strong><br>"
        '<code>DA/</code> · <code>SA/</code> · <code>HV/</code></li>',
    ]
    title_by_stem = {stem: title for stem, title in STEM_MAP.values()}
    listed: set[str] = set()
    for name, (stem, title) in STEM_MAP.items():
        src = args.src_dir / name
        if not src.exists():
            print("SKIP missing", src)
            continue
        out = convert_one(src, args.out_dir)
        print("OK", out.name)
        sub = finished_manuals_subdir(stem)
        listed.add(stem)
        index_rows.append(
            f"<li><strong>{html.escape(title)}</strong><br>"
            f'<a href="{sub}/{stem}.html">HTML (сетка 12 / A4)</a></li>',
        )

    for fam_dir, stem, _html_path in iter_finished_manual_html(args.out_dir):
        if stem in listed:
            continue
        label = title_by_stem.get(stem, stem)
        index_rows.append(
            f"<li><strong>{html.escape(label)}</strong><br>"
            f'<a href="{fam_dir}/{stem}.html">HTML (сетка 12 / A4)</a></li>',
        )

    (args.out_dir / "index.html").write_text(
        "<!DOCTYPE html><html lang='ru'><head><meta charset='utf-8'>"
        "<title>Инструкции SKU (RU)</title>"
        "<style>body{font-family:system-ui;max-width:820px;margin:2rem auto;"
        "padding:0 1rem}li{margin:.65rem 0}</style></head><body>"
        "<h1>Руководства по эксплуатации (RU)</h1>"
        "<p>Адаптивная вёрстка на 12 колонках, 2 листа A4 альбом. "
        "Готовые руководства — в <code>DA/</code>, <code>SA/</code>, "
        "<code>HV/</code>. "
        "Текст и схемы — по <code>_инструкции-pdf/</code>; термины — "
        "<code>docs/tech-copy-belimo-ru.md</code>. "
        "Новые переводы с раздельными PDF 24/230 В — шаблоны "
        "<a href='template-v24.html'>V24</a> / "
        "<a href='template-v230.html'>V230</a>.</p>"
        f"<ol>{''.join(index_rows)}</ol></body></html>\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
